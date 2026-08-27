# -*- coding: utf-8 -*-
"""P7 재작성: 레버리지가 항상 불리하진 않다 — 추세장에서는 유리
- 좌: 추세 상승장(2026.03.31~05.29) 실제 +207% vs 단순 2배 +159% (+48%p 더 번다)
- 우: 추세 하락장(2022.01.04~10.04) 실제 -48.0% vs 단순 2배 -52.3% (+4.3%p 덜 잃는다)
데이터: ETF_데이터_pivot (표기는 ***** 마스킹)
"""
import sys, io, os, platform
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
import numpy as np, pandas as pd
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
plt.rcParams['font.family']=('Malgun Gothic' if platform.system()=='Windows' else 'AppleGothic')
plt.rcParams['axes.unicode_minus']=False
NAVY='#043B72'; ORANGE='#F58220'; GRAY='#84888B'
os.chdir(os.path.join("C:/Users/gamer38/Documents/Claude/Projects","변동성펀드","제안서"))
df=pd.read_parquet("C:/Users/gamer38/Documents/Claude/Projects/ETF WEB/ETF_데이터_pivot.parquet")
def px(nm,col='수정주가(원)'):
    k=df[df['종목명']==nm].dropna(subset=[col]).sort_values('날짜')
    s=pd.Series(k[col].astype(float).values,index=pd.to_datetime(k['날짜']).values)
    i=pd.to_datetime(s.index); s=s[i.dayofweek<5]
    chg=s.pct_change().fillna(1.0); return s[chg!=0]
k1=px('KODEX 200'); k2=px('KODEX 레버리지')

def style(ax):
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.grid(axis='y',color='#e6ebf1',lw=0.7); ax.set_axisbelow(True)
    ax.tick_params(labelsize=8.5)

wins=[("추세 상승장 : 단순 2배보다 더 번다","2026-03-31","2026-05-29"),
      ("추세 하락장 : 단순 2배보다 덜 잃는다","2022-01-04","2022-10-04")]
fig,axes=plt.subplots(1,2,figsize=(10.4,3.4),dpi=150)
for ax,(tt,a,b) in zip(axes,wins):
    w1=k1[(k1.index>=pd.Timestamp(a))&(k1.index<=pd.Timestamp(b))]
    w2=k2[(k2.index>=pd.Timestamp(a))&(k2.index<=pd.Timestamp(b))]
    rb=w1/w1.iloc[0]; rl=w2/w2.iloc[0]
    simple=1.0+(rb-1.0)*2.0
    ax.plot(rl.index,rl.values*100,color=ORANGE,lw=2.0,label='레버리지 실제 (복리)')
    ax.plot(simple.index,simple.values*100,color=GRAY,lw=1.6,ls='--',label='단순 2배 (가상)')
    ax.plot(rb.index,rb.values*100,color=NAVY,lw=1.4,label='KOSPI200')
    span=(max(rl.max(),simple.max(),rb.max())-min(rl.min(),simple.min(),rb.min()))*100
    labs=sorted([(float(rl.iloc[-1])*100,ORANGE),(float(simple.iloc[-1])*100,GRAY),(float(rb.iloc[-1])*100,NAVY)],key=lambda t:-t[0])
    pos=[]
    for v,_ in labs: pos.append(v if not pos else min(v,pos[-1]-span*0.09))
    for (v,c),py in zip(labs,pos):
        ax.annotate(f"{v-100:+.1f}%",xy=(rl.index[-1],py),xytext=(5,-3),
                    textcoords='offset points',fontsize=9.5,fontweight='bold',color=c)
    adv=(float(rl.iloc[-1])-float(simple.iloc[-1]))*100
    lab=f"복리 효과 : {adv:+.0f}%p 더 번다" if adv>0 and float(rl.iloc[-1])>1 else f"복리 효과 : {adv:+.1f}%p 덜 잃는다"
    g=float(rb.iloc[-1]-1)
    ax.annotate(lab,xy=(0.04,0.90) if g<0 else (0.04,0.90),xycoords='axes fraction',
                fontsize=10.5,fontweight='bold',color='#b45309')
    ax.set_title(f"{tt}\n({a[:7].replace('-','.')} ~ {b[:7].replace('-','.')} · KOSPI200 {g*100:+.0f}%)",
                 fontsize=10.5,fontweight='bold',color='#222')
    ax.legend(fontsize=7.5,loc=('upper left' if g>0 else 'lower left'),frameon=False,
              bbox_to_anchor=((0.0,0.88) if g>0 else (0.0,0.02)))
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%y.%m'))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=1 if g>0 else 2)); style(ax)
    ax.set_xlim(rl.index[0],rl.index[-1]+pd.Timedelta(days=10 if g>0 else 24))
    print(f"[{tt}] 지수 {g*100:+.1f}% · 단순2배 {(float(simple.iloc[-1])-1)*100:+.1f}% · 실제 {(float(rl.iloc[-1])-1)*100:+.1f}% ({adv:+.1f}%p)")
plt.tight_layout()
plt.savefig(os.path.join('img','lev15_trend2.png')); plt.close()
print("차트 저장: lev15_trend2.png")

# ---- P7 재작성 ----
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
DST="Tech_TOP10_스마트액티브레버리지15_v1_20260812.pptx"
prs=Presentation(DST)
DARK=RGBColor(0x22,0x22,0x22); GRAY_=RGBColor(0x84,0x88,0x8B); ORANGE_=RGBColor(0xE8,0x72,0x0C)
FT_B="KoPub돋움체_Pro Bold"; FT_M="KoPub돋움체_Pro Medium"
SW=prs.slide_width/914400; SH=prs.slide_height/914400
def tbox(s,L,T,W,H,text,size,bold,color,font=None,sp=1.15):
    b=s.shapes.add_textbox(Inches(L),Inches(T),Inches(W),Inches(H))
    tf=b.text_frame; tf.word_wrap=True
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
        if font: r.font.name=font
        p.line_spacing=sp
    return b
s=prs.slides[6]
for sh in list(s.shapes): s.shapes._spTree.remove(sh._element)
tbox(s,0.54,0.42,SW-1.1,0.45,"2. 레버리지 ETF의 특징 ② : 추세장에서는 오히려 유리하다",21,True,DARK,FT_B)
tbox(s,0.55,0.98,SW-1.1,0.6,
     "레버리지가 항상 불리한 것은 아닙니다. 한 방향 추세에서는 일별 복리가 유리하게 작동합니다 — "
     "상승 추세에서는 단순 2배보다 더 벌고, 하락 추세에서는 단순 2배보다 덜 잃습니다 (KOSPI200 실측).",13.5,None,DARK)
s.shapes.add_picture(os.path.join('img','lev15_trend2.png'),Inches(0.45),Inches(2.05),width=Inches(9.95))
tbox(s,0.8,5.85,SW-1.7,0.5,"왕복 구간에는 불리하고, 추세 구간에는 유리합니다 — 방향이 아니라 '경로'가 성과를 결정합니다.",14,True,ORANGE_,FT_B)
tbox(s,0.55,SH-0.62,SW-1.1,0.4,
     "※ ***** 200·***** 레버리지 수정주가 실측 (상승 사례 2026.03.31~05.29, 하락 사례 2022.01.04~10.04) · 단순 2배 = 기초 누적수익률×2 (가상) · 데이터: ETF데이터 · 상기 자료는 이해를 돕기 위한 예시입니다.",8,False,GRAY_,FT_M)
prs.save(DST)
print("P7 저장 완료")
