# -*- coding: utf-8 -*-
"""재구성 2차: 목차 겹침 수정 · 펀드정보/투자위험 내용 채움 · 꼬리말 Tech TOP10"""
import sys, io, os
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
os.chdir(os.path.join("C:/Users/gamer38/Documents/Claude/Projects","변동성펀드","제안서"))
DST="Tech_TOP10_스마트액티브레버리지15_v1_20260812.pptx"
prs=Presentation(DST)
DARK=RGBColor(0x22,0x22,0x22); NAVY=RGBColor(0x04,0x3B,0x72); GRAY=RGBColor(0x84,0x88,0x8B)
WHITE=RGBColor(0xFF,0xFF,0xFF)
FT_B="KoPub돋움체_Pro Bold"; FT_M="KoPub돋움체_Pro Medium"
SW=prs.slide_width/914400; SH=prs.slide_height/914400

def tbox(s,L,T,W,H,text,size,bold,color,font=None,sp=1.15):
    b=s.shapes.add_textbox(Inches(L),Inches(T),Inches(W),Inches(H))
    tf=b.text_frame; tf.word_wrap=True
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
        if font: r.font.name=font
        p.line_spacing=sp
    return b
def para_replace(p,a,b):
    runs=p.runs; full=''.join(r.text for r in runs)
    if a not in full: return 0
    st=full.index(a); en=st+len(a); pos=0; first=True
    for r in runs:
        ln=len(r.text); s0,e0=pos,pos+ln
        if e0<=st or s0>=en: pos=e0; continue
        pre=r.text[:max(st-s0,0)]; post=r.text[max(en-s0,0):] if en-s0<ln else ''
        if first: r.text=pre+b+post; first=False
        else: r.text=pre+post
        pos=e0
    return 1

# ---- 1) 목차(2p): 내가 넣은 큰 '목차' 제목만 제거 (레이아웃 제목 유지) ----
s2=prs.slides[1]
for sh in list(s2.shapes):
    if getattr(sh,'has_text_frame',False) and sh.has_text_frame and sh.text_frame.text.strip()=="목차":
        s2.shapes._spTree.remove(sh._element); print("목차 중복 제목 제거")

# ---- 2) 18p 펀드 정보: 표 추가 (제목은 기존 유지) ----
s18=prs.slides[17]
rows=[("펀드명 (가칭)","미래에셋 Tech TOP10 스마트 액티브 레버리지 1.5 증권투자신탁 (주식-파생형)"),
      ("투자대상","반도체 TOP10 지수 관련 ETF(TIGER 반도체TOP10 등) 및 장내 파생상품"),
      ("목표 노출","평균 1.5배 (일간 배수를 고정하지 않는 스마트 액티브 방식)"),
      ("운용방식","규칙 기반 리밸런싱 + 운용역 판단 병행 (액티브) · QPMS Vol Trading 시스템 활용"),
      ("비교지수 (BM)","반도체 TOP10 지수 일간수익률 ×1.5 (매일 재조정 가정)"),
      ("위험등급","1등급 (매우 높은 위험) — 레버리지형"),
      ("운용","미래에셋자산운용 AI금융공학운용부문"),
      ("보수·판매","추후 확정")]
tb=s18.shapes.add_table(len(rows),2,Inches(0.8),Inches(1.5),Inches(SW-1.6),Inches(4.6)).table
tb.columns[0].width=Inches(2.2); tb.columns[1].width=Inches(SW-1.6-2.2)
for i,(k,v) in enumerate(rows):
    c0,c1=tb.cell(i,0),tb.cell(i,1)
    c0.text=k; c1.text=v
    for c,bold,col in [(c0,True,WHITE),(c1,False,DARK)]:
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size=Pt(11.5); r.font.bold=bold; r.font.color.rgb=col
                r.font.name=FT_B if bold else FT_M
    c0.fill.solid(); c0.fill.fore_color.rgb=NAVY
    c1.fill.solid(); c1.fill.fore_color.rgb=RGBColor(0xF8,0xFA,0xFC) if i%2 else WHITE
tbox(s18,0.8,6.35,SW-1.6,0.4,"※ 상기 내용은 제안 단계의 초안이며, 신탁계약·투자설명서 확정 과정에서 변경될 수 있습니다.",9,False,GRAY,FT_M)
print("18p 펀드 정보 표 작성")

# ---- 3) 19p 투자 위험: 내용 채움 (제목은 기존 유지) ----
s19=prs.slides[18]
risks=[("원금 손실 위험","이 펀드는 예금자보호법에 따라 보호되지 않으며, 투자원금 전액 손실이 발생할 수 있습니다. 투자 결과는 전부 투자자에게 귀속됩니다."),
       ("레버리지 위험","평균 1.5배 노출로 기초지수 하락 시 손실이 지수보다 확대됩니다. 백테스트 기간 최대낙폭은 약 -69%로 매우 깊습니다."),
       ("경로 의존(변동성 잠식) 위험","평균 배수 유지 방식은 일일 재조정 비용을 줄일 뿐, 왕복 구간에서 발생하는 레버리지 고유의 손실 구조 자체를 없애지 못합니다."),
       ("집중투자 위험","반도체 10종목에 집중 투자하며 삼성전자·SK하이닉스 비중이 높아, 개별기업·산업 사이클 위험에 크게 노출됩니다."),
       ("파생상품·추적편차 위험","장내 파생상품 활용에 따른 베이시스·롤오버 비용이 발생할 수 있고, 평균 1.5배 추종 특성상 단기 성과는 일간 1.5배와 다를 수 있습니다."),
       ("유동성·시장 위험","시장 급변동 시 재조정 매매의 체결 가격이 불리할 수 있으며, 거래 부진 시 환매 대응이 지연될 수 있습니다.")]
y=1.45
for h,b in risks:
    tbox(s19,0.8,y,SW-1.7,0.3,h,13,True,NAVY,FT_B)
    tbox(s19,0.8,y+0.34,SW-1.7,0.55,b,10.5,None,DARK,sp=1.2)
    y+=0.92
tbox(s19,0.8,y+0.05,SW-1.7,0.4,"※ 자세한 위험은 투자설명서를 반드시 확인하시기 바랍니다.",9,False,GRAY,FT_M)
print("19p 투자 위험 작성")

# ---- 4) 꼬리말: IT TOP10 → Tech TOP10 (마스터·레이아웃·슬라이드 전체) ----
n=0
def sweep(shapes):
    global n
    for sh in shapes:
        if getattr(sh,'has_text_frame',False) and sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                n+=para_replace(p,"IT TOP10","Tech TOP10")
for m in prs.slide_masters:
    sweep(m.shapes)
    for lo in m.slide_layouts: sweep(lo.shapes)
for s in prs.slides: sweep(s.shapes)
print(f"IT TOP10 → Tech TOP10 교체 {n}건")

prs.save(DST)
print("저장 완료")
