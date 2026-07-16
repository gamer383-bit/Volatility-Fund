# -*- coding: utf-8 -*-
"""AITop2 v1 제안서에 '지수대별 실제 편입비' 슬라이드 삽입 (새 10페이지, HOW 다음)
- 성장형(풋매도+KO풋+콜스프레드 K110/150)·안정형(풋매도) 두 테이블
- 행=지수 60~140(10pt), 열=잔존만기 12/9/6/3/1개월, σ60%, r3%, q2%, 편입비 캡 0~100%
"""
import sys, io, os, math
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
import numpy as np
from scipy.special import ndtr as N
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

R,Q,SIG=0.03,0.02,0.60
KPUT,KKO,H,K1,K2=100.,100.,60.,110.,150.
NAVY=RGBColor(0x04,0x3B,0x72); ORANGE=RGBColor(0xF5,0x82,0x20); GRAY=RGBColor(0x84,0x88,0x8B)
DGRAY=RGBColor(0x3A,0x3A,0x3A); LIGHT=RGBColor(0xEE,0xF2,0xF7); WHITE=RGBColor(0xFF,0xFF,0xFF)
FT="맑은 고딕"

def bs_delta(typ,S,K,T,sig=SIG):
    sq=sig*math.sqrt(T); d1=(math.log(S/K)+(R-Q+0.5*sig*sig)*T)/sq
    eqt=math.exp(-Q*T)
    return eqt*N(d1) if typ=='c' else eqt*(N(d1)-1)
def bs_price(typ,S,K,T,sig=SIG):
    sq=sig*math.sqrt(T); d1=(math.log(S/K)+(R-Q+0.5*sig*sig)*T)/sq; d2=d1-sq
    eqt=math.exp(-Q*T); ert=math.exp(-R*T)
    if typ=='c': return S*eqt*N(d1)-K*ert*N(d2)
    return K*ert*N(-d2)-S*eqt*N(-d1)
def din_put(S,K,Hb,T,sig=SIG):
    b=R-Q; sq=sig*math.sqrt(T); mu=(b-0.5*sig*sig)/(sig*sig); phi=-1.; eta=1.
    x2=math.log(S/Hb)/sq+(1+mu)*sq; y1=math.log(Hb*Hb/(S*K))/sq+(1+mu)*sq; y2=math.log(Hb/S)/sq+(1+mu)*sq
    ebr=math.exp((b-R)*T); ert=math.exp(-R*T); pw1=(Hb/S)**(2*(mu+1)); pw2=(Hb/S)**(2*mu)
    B=phi*S*ebr*N(phi*x2)-phi*K*ert*N(phi*x2-phi*sq)
    C=phi*S*ebr*pw1*N(eta*y1)-phi*K*ert*pw2*N(eta*y1-eta*sq)
    D=phi*S*ebr*pw1*N(eta*y2)-phi*K*ert*pw2*N(eta*y2-eta*sq)
    return B-C+D
def dop(S,K,Hb,T):
    if S<=Hb: return 0.0
    return max(bs_price('p',S,K,T)-din_put(S,K,Hb,T),0.0)
def ko_delta(S,K,Hb,T):
    if S<=Hb: return 0.0
    h=max(S*1e-4,1e-6)
    return (dop(S+h,K,Hb,T)-dop(max(S-h,Hb+1e-9),K,Hb,T))/(2*h)

def w_stable(S,tau):   # 안정형: 풋매도 복제
    return min(max(-bs_delta('p',S,KPUT,tau),0.0),1.0)*100
def w_growth(S,tau):   # 성장형: 풋매도 + KO풋 + 콜스프레드
    d=-bs_delta('p',S,KPUT,tau)+ko_delta(S,KKO,H,tau)+(bs_delta('c',S,K1,tau)-bs_delta('c',S,K2,tau))
    return min(max(d,0.0),1.0)*100

LEVELS=[140,130,120,110,100,90,80,70,60]
TAUS=[(12,1.0),(9,0.75),(6,0.5),(3,0.25),(1,1/12)]

def build_rows(fn):
    rows=[["지수 \\ 잔존"]+[f"{m}개월" for m,_ in TAUS]]
    for lv in LEVELS:
        r=[f"{lv}"]
        for _,tau in TAUS:
            if lv<=H and fn is w_growth: r.append("낙아웃")     # 배리어 도달
            else: r.append(f"{fn(float(lv),tau):.0f}%")
        rows.append(r)
    return rows

g_rows=build_rows(w_growth); s_rows=build_rows(w_stable)
print("=== 성장형 ===");  [print(" | ".join(r)) for r in g_rows]
print("=== 안정형 ===");  [print(" | ".join(r)) for r in s_rows]

# ---------- 슬라이드 삽입 ----------
PATH="C:/Users/gamer38/Documents/Claude/Projects/변동성펀드/제안서/변동성하베스트펀드_AITop2_v1_20260715.pptx"
prs=Presentation(PATH)
s=prs.slides.add_slide(prs.slide_layouts[6])

def rect(x,y,w,h,color):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    r.fill.solid(); r.fill.fore_color.rgb=color; r.line.fill.background(); r.shadow.inherit=False
    return r
def txt(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=sp
        for (t,sz,col,b) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=col; r.font.name=FT
    return tb
def table(x,y,w,rows,hdr_fill,fs=10):
    nr,nc=len(rows),len(rows[0]); rh=Inches(0.33)
    gt=s.shapes.add_table(nr,nc,x,y,w,rh*nr).table
    cw=[Inches(0.95)]+[ (w-Inches(0.95))//(nc-1) ]*(nc-1)
    for j,c in enumerate(cw): gt.columns[j].width=c
    for i,row in enumerate(rows):
        gt.rows[i].height=rh
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.margin_left=Pt(3); c.margin_right=Pt(3); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=str(val)
            r.font.name=FT; r.font.size=Pt(fs); p.alignment=PP_ALIGN.CENTER
            if i==0:
                c.fill.solid(); c.fill.fore_color.rgb=hdr_fill; r.font.color.rgb=WHITE; r.font.bold=True
            else:
                c.fill.solid(); c.fill.fore_color.rgb=(LIGHT if i%2==0 else WHITE)
                r.font.color.rgb=DGRAY; r.font.bold=(j==0 or row[0]=="100")

# 배경 흰색
bg=rect(0,0,prs.slide_width,prs.slide_height,WHITE)
s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2,bg._element)
# 헤더 (덱 스타일)
rect(Inches(0.5),Inches(0.55),Inches(0.14),Inches(0.62),ORANGE)
txt(Inches(0.75),Inches(0.5),Inches(9),Inches(0.35),[[("08  HOW",12,ORANGE,True)]])
txt(Inches(0.75),Inches(0.8),Inches(9.5),Inches(0.55),[[("지수대별 실제 편입비 — 만기까지 남은 시간 활용",23,NAVY,True)]])
txt(Inches(0.75),Inches(1.42),Inches(9.5),Inches(0.4),[[("리밸런싱일 지수=100 기준 · 편입비는 잔존만기가 짧아질수록 가팔라짐 · 최대 편입비 100%",12.5,DGRAY,True)]])
# 좌: 성장형
txt(Inches(0.6),Inches(1.95),Inches(4.7),Inches(0.35),[[("성장변동성펀드",14,ORANGE,True),("  (변동성 매매 + 상승 참여)",10.5,GRAY,False)]])
table(Inches(0.6),Inches(2.32),Inches(4.75),g_rows,ORANGE,fs=9.5)
# 우: 안정형
txt(Inches(5.6),Inches(1.95),Inches(4.7),Inches(0.35),[[("안정변동성펀드",14,NAVY,True),("  (안정 변동성 매매)",10.5,GRAY,False)]])
table(Inches(5.6),Inches(2.32),Inches(4.75),s_rows,NAVY,fs=9.5)
# 각주
txt(Inches(0.6),Inches(5.85),Inches(9.6),Inches(0.7),[
    [("· 변동성 60%·금리 3%·배당 2% 가정, 지수 10포인트 단위(60~140). 성장형은 지수 60 이하 도달 시 '낙아웃'(하방 완충 소멸) 후 주식형과 유사하게 운용.",9.5,GRAY,False)],
    [("· 표의 편입비는 모형상 목표치이며, 실제 운용에서는 자체 드리프트·거래비용을 감안해 소폭 다를 수 있습니다.",9.5,GRAY,False)]],sp=1.25)
txt(Inches(0.5),Inches(7.12),Inches(8.6),Inches(0.32),[[("※ 상기 자료는 이해를 돕기 위한 예시이며 실제 투자내용을 의미하지 않습니다. 시뮬레이션은 가정에 따른 것으로 실제와 다를 수 있습니다.",7,GRAY,False)]])

# 위치 이동: HOW(현재 9번째, idx8) 바로 뒤 → 10페이지(idx9)
xml=prs.slides._sldIdLst; ids=list(xml)
xml.remove(ids[-1]); xml.insert(9,ids[-1])

try:
    prs.save(PATH); print("\n[OK] 새 10페이지 삽입 완료:", PATH)
except PermissionError:
    print("\n[잠김] PowerPoint에 열려 있어 저장 불가 — 파일을 닫고 다시 요청해주세요.")
