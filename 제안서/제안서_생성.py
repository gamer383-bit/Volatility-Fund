# -*- coding: utf-8 -*-
"""변동성 하베스트 펀드 제안서 2종(KOSPI200 / AI Top2) 생성 — 미래에셋 형식 초안
용어 방침: '변동성 매도/프리미엄 수취' 등 금융공학 용어 대신
'변동성 매매를 통해 수익을 추구' 톤으로 통일 (금융공학 색 최소화)"""
import math, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 색상/폰트 ----------
NAVY=RGBColor(0x04,0x3B,0x72); ORANGE=RGBColor(0xF5,0x82,0x20); CYAN=RGBColor(0x00,0xA9,0xCE)
GRAY=RGBColor(0x84,0x88,0x8B); DGRAY=RGBColor(0x3A,0x3A,0x3A); LIGHT=RGBColor(0xEE,0xF2,0xF7)
WHITE=RGBColor(0xFF,0xFF,0xFF); LORANGE=RGBColor(0xF0,0xB2,0x6B); RED=RGBColor(0xD0,0x2F,0x00)
GREEN=RGBColor(0x1F,0x9E,0x6E)
FT="맑은 고딕"; FTB="맑은 고딕"
SW,SH=Inches(10.8),Inches(7.5)
IMG=os.path.join(os.path.dirname(os.path.abspath(__file__)),"img")   # 차트 PNG 폴더

# ---------- 옵션 계산 ----------
def N(x): return 0.5*(1+math.erf(x/math.sqrt(2)))
def bs(t,S,K,T,sig,r,q):
    if T<=0: return max(S-K,0) if t=='c' else max(K-S,0)
    sq=sig*math.sqrt(T); d1=(math.log(S/K)+(r-q+0.5*sig*sig)*T)/sq; d2=d1-sq
    if t=='c': return S*math.exp(-q*T)*N(d1)-K*math.exp(-r*T)*N(d2)
    return K*math.exp(-r*T)*N(-d2)-S*math.exp(-q*T)*N(-d1)
def scenario(sig):
    S0=100;T=1;r=0.03;q=0.02;putK=100;k1=110;k2=140
    pput=bs('p',S0,putK,T,sig,r,q); psp=bs('c',S0,k1,T,sig,r,q)-bs('c',S0,k2,T,sig,r,q); g=math.exp(r*T)
    rows=[]
    for m in [60,70,80,90,100,110,120,130,140,150]:
        stable=pput*g-max(putK-m,0)
        active=stable+((max(m-k1,0)-max(m-k2,0))-psp*g)
        rows.append((m,stable,active))
    return pput,psp,rows

# ---------- 그리기 헬퍼 ----------
def bg(slide,color):
    r=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=color; r.line.fill.background()
    r.shadow.inherit=False
    slide.shapes._spTree.remove(r._element); slide.shapes._spTree.insert(2,r._element)
    return r
def rect(slide,x,y,w,h,color,line=None):
    r=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    r.fill.solid(); r.fill.fore_color.rgb=color
    if line is None: r.line.fill.background()
    else: r.line.color.rgb=line; r.line.width=Pt(0.75)
    r.shadow.inherit=False; return r
def txt(slide,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0,wrap=True):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=wrap
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    tf.vertical_anchor=anchor
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=sp
        for (t,sz,col,b) in para:
            r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.bold=b; f.color.rgb=col; f.name=FT
    return tb
def bullets(slide,x,y,w,h,items,sz=13,col=DGRAY,gap=1.15):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(mark,t,c,b) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=gap; p.space_after=Pt(3)
        r=p.add_run(); r.text=(mark+" " if mark else "")+t
        r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=b; r.font.name=FT
    return tb
def footer(slide,pg,note="※ 상기 자료는 이해를 돕기 위한 예시이며 실제 투자내용을 의미하지 않습니다. 시뮬레이션은 가정에 따른 것으로 실제와 다를 수 있습니다."):
    txt(slide,Inches(0.5),Inches(7.12),Inches(8.6),Inches(0.32),[[(note,7,GRAY,False)]])
    txt(slide,Inches(9.9),Inches(7.12),Inches(0.6),Inches(0.32),[[(str(pg),9,GRAY,True)]],align=PP_ALIGN.RIGHT)
def header(slide,kicker,title):
    rect(slide,Inches(0.5),Inches(0.55),Inches(0.14),Inches(0.62),ORANGE)
    txt(slide,Inches(0.75),Inches(0.5),Inches(9),Inches(0.35),[[(kicker,12,ORANGE,True)]])
    txt(slide,Inches(0.75),Inches(0.8),Inches(9.5),Inches(0.55),[[(title,23,NAVY,True)]])

def make_table(slide,x,y,w,rows,colw,header_fill=NAVY,fs=11,hrow=True,rh=Inches(0.34)):
    nr=len(rows); nc=len(rows[0])
    gt=slide.shapes.add_table(nr,nc,x,y,w,rh*nr).table
    gt.first_row=False; gt.horz_banding=False
    # remove default style banding by setting col widths
    for j,cw in enumerate(colw): gt.columns[j].width=cw
    for i,row in enumerate(rows):
        gt.rows[i].height=rh
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.margin_left=Pt(5); c.margin_right=Pt(5); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            tf=c.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
            r=p.add_run(); r.text=str(val); r.font.name=FT; r.font.size=Pt(fs)
            if hrow and i==0:
                c.fill.solid(); c.fill.fore_color.rgb=header_fill
                r.font.color.rgb=WHITE; r.font.bold=True; p.alignment=PP_ALIGN.CENTER
            else:
                c.fill.solid(); c.fill.fore_color.rgb=(LIGHT if i%2==0 else WHITE)
                r.font.color.rgb=DGRAY; r.font.bold=(j==0)
                p.alignment=(PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
    return gt

def pic(slide,path,x,y,w=None,h=None):
    kw={}
    if w is not None: kw['width']=w
    if h is not None: kw['height']=h
    return slide.shapes.add_picture(path,x,y,**kw)

# =====================================================================
def build(path, name, under, sig, is_index, volimg, tag):
    pput,psp,rows=scenario(sig)
    valu_txt=("KOSPI 저평가 — 글로벌 대비 낮은 PER·PBR" if is_index
              else "삼성전자·SK하이닉스 저평가 — 실적 대비 낮은 밸류에이션")
    prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
    blank=prs.slide_layouts[6]
    S=lambda: prs.slides.add_slide(blank)

    # ---- 1) 표지 ----
    s=S(); bg(s,NAVY)
    rect(s,0,Inches(2.55),SW,Inches(0.06),ORANGE)
    txt(s,Inches(0.9),Inches(1.4),Inches(9),Inches(0.5),[[("미래에셋자산운용  AI금융공학운용부문",14,CYAN,True)]])
    txt(s,Inches(0.9),Inches(2.75),Inches(9.6),Inches(1.4),[
        [("변동성 하베스트 펀드",40,WHITE,True)],
        [(name,26,LORANGE,True)]],sp=1.05)
    txt(s,Inches(0.9),Inches(5.2),Inches(9),Inches(0.4),[[("역대급으로 높아진 변동성을 수익으로 수확(Harvest)하는 변동성 매매 전략",15,RGBColor(0xC9,0xD6,0xE5),False)]])
    txt(s,Inches(0.9),Inches(6.4),Inches(9),Inches(0.4),[[("2026.07  |  내부 검토용 초안",12,GRAY,False)]])
    txt(s,Inches(0.9),Inches(6.75),Inches(9.4),Inches(0.4),[[("미래에셋자산운용 컴플라이언스 검토 전 초안 — 대외 배포 금지",10,LORANGE,False)]])

    # ---- 2) 상품 개요 ----
    s=S(); bg(s,WHITE); header(s,"01  PRODUCT","상품 개요")
    ov=[["구  분","내  용"],
        ["명  칭",f"미래에셋 변동성 하베스트 펀드 — {name} (가칭)"],
        ["개  요","기초자산 편입비를 시장 상황에 따라 조절하는 변동성 매매를 통해, 역대급으로 높아진 변동성을 수익으로 전환하는 전략"],
        ["기초자산",under],
        ["가정 변동성",f"연 {int(sig*100)}% (최근 실현변동성 대비 보수적 가정)"],
        ["운용 전략","변동성 매매(저가매수·고가매도)로 수익 추구 + 상승 방향 참여 · 최대 편입비 180% 제한"],
        ["유형 구분","안정변동성펀드(안정 변동성 매매) / 성장변동성펀드(변동성 매매 + 상승 참여)"],
        ["상품 유형","국내 주식형 파생 (사모 · 일임)"],
        ["투자 등급","1등급 (매우 높은 위험) · 적극투자형"],
        ["만  기","1년 (예시) · 일별 리밸런싱"],
        ["최저가입금액","3천만원"]]
    make_table(s,Inches(0.75),Inches(1.5),Inches(9.3),ov,[Inches(1.7),Inches(7.6)],fs=11,rh=Inches(0.44))
    txt(s,Inches(0.75),Inches(6.55),Inches(9.4),Inches(0.5),[
        [("※ 원금 손실(0~100%)이 발생할 수 있으며 손실은 투자자에게 귀속됩니다. 예금자보호법 적용 대상이 아닙니다.",8.5,GRAY,False)]])
    footer(s,2)

    # ---- 왜 이 펀드인가 ① 역대급 변동성 ----
    s=S(); bg(s,WHITE); header(s,"02  WHY NOW","왜 지금 이 펀드인가 ① — 역대급 변동성")
    txt(s,Inches(0.75),Inches(1.05),Inches(9.5),Inches(0.4),[[("최근 한국시장 변동성이 5년 내 최고 수준(99.7%ile). 변동성이 클수록 변동성 매매의 수익 기회가 커집니다.",13,DGRAY,True)]])
    pic(s,volimg,Inches(0.55),Inches(1.5),w=Inches(9.7))
    rect(s,Inches(0.75),Inches(6.0),Inches(9.3),Inches(0.72),RGBColor(0xFD,0xF1,0xE6))
    txt(s,Inches(0.95),Inches(6.0),Inches(8.9),Inches(0.72),[[("→ 지금은 변동성을 '피할' 때가 아니라 '수확'하기 가장 유리한 국면",13,ORANGE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s,3)

    # ---- 왜 ② 당분간 제한적 방향성 ----
    s=S(); bg(s,WHITE); header(s,"03  WHY NOW","왜 지금 이 펀드인가 ② — 당분간 제한적 방향성")
    txt(s,Inches(0.75),Inches(1.32),Inches(9.5),Inches(0.4),[[("상방엔 매물벽, 하방엔 저평가 — 한쪽으로 크게 쏠린 방향성이 나오기 어려운 국면입니다.",13,DGRAY,True)]])
    txt(s,Inches(0.75),Inches(1.82),Inches(9.5),Inches(0.32),[[("① 상방을 누르는 수급 요인",13,NAVY,True)]])
    cards=[("국민연금 매도 대기","목표비중 초과·국내주식 축소 기조로, 지수가 오를수록 대기 매도 물량이 상단에서 출회."),
           ("외국인 리밸런싱","반도체 쏠림에 따른 지수 편입비 조정·차익 실현으로 상승 국면에서 매도 우위가 반복."),
           ("개인 순응매도","단기 급등에 올라탄 과도한 매수 포지션이 반등 시마다 차익실현(상승 시 매도 대응)으로 전환.")]
    for i,(t,d) in enumerate(cards):
        cx=Inches(0.75+i*3.15)
        rect(s,cx,Inches(2.22),Inches(2.95),Inches(2.5),LIGHT)
        rect(s,cx,Inches(2.22),Inches(2.95),Inches(0.5),NAVY)
        txt(s,cx+Inches(0.15),Inches(2.27),Inches(2.7),Inches(0.4),[[(t,12.5,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
        txt(s,cx+Inches(0.18),Inches(2.88),Inches(2.6),Inches(1.75),[[(d,11.5,DGRAY,False)]],sp=1.22)
    rect(s,Inches(0.75),Inches(4.95),Inches(9.3),Inches(0.9),RGBColor(0xEC,0xF3,0xEA))
    txt(s,Inches(0.95),Inches(5.02),Inches(8.9),Inches(0.78),[
        [("② 하방을 받치는 밸류에이션: ",12.5,GREEN,True),(valu_txt,12.5,DGRAY,True)],
        [("   ※ 세부 밸류에이션 수치는 별도 보고서 반영 예정",9.5,GRAY,False)]],sp=1.15)
    rect(s,Inches(0.75),Inches(6.02),Inches(9.3),Inches(0.7),RGBColor(0xFD,0xF1,0xE6))
    txt(s,Inches(0.95),Inches(6.02),Inches(8.9),Inches(0.7),[[("→ 상방 제한 + 하방 지지 = 좁은 등락 반복 → 방향성 베팅보다 변동성 매매가 유리",13,ORANGE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s,4)

    # ---- 왜 ③ 레버리지가 키운 변동성 ----
    s=S(); bg(s,WHITE); header(s,"04  WHY NOW","왜 지금 이 펀드인가 ③ — 레버리지가 키운 변동성")
    txt(s,Inches(0.75),Inches(1.05),Inches(9.5),Inches(0.4),[[("단일종목 레버리지 ETF 급증이 삼성전자·SK하이닉스와 지수 변동성을 구조적으로 키우고 있습니다.",13,DGRAY,True)]])
    pic(s,os.path.join(IMG,"lev_aum.png"),Inches(0.55),Inches(1.5),w=Inches(4.55))
    pic(s,os.path.join(IMG,"idx_weight.png"),Inches(5.5),Inches(1.5),w=Inches(4.55))
    rect(s,Inches(0.75),Inches(4.62),Inches(9.3),Inches(2.08),RGBColor(0xF2,0xF5,0xF8))
    txt(s,Inches(0.95),Inches(4.72),Inches(8.9),Inches(0.35),[[("변동성이 변동성을 부르는 되먹임 — 두 종목이 KOSPI200의 ",13,NAVY,True),("51.5%",13,ORANGE,True),(" (연초 38.7%→5월)",12,DGRAY,False)]])
    loop=["① 지수·종목 흔들림","② 지수의 절반인\n두 종목 급등락","③ 레버리지 종가\n순응매매(±10%→2.6조)","④ 지수 변동성\n재확대"]
    for i,tt in enumerate(loop):
        bx=Inches(0.95+i*2.30)
        rect(s,bx,Inches(5.35),Inches(2.02),Inches(1.05),WHITE,line=NAVY)
        txt(s,bx+Inches(0.06),Inches(5.35),Inches(1.9),Inches(1.05),[[(tt,10.5,NAVY,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=1.05)
        if i<3: txt(s,Inches(2.97+i*2.30),Inches(5.55),Inches(0.28),Inches(0.6),[[("→",16,ORANGE,True)]],align=PP_ALIGN.CENTER)
    footer(s,5)

    # ---- 잠식 역이용: 레버리지가 잃는 변동성을 우리는 수익으로 ----
    s=S(); bg(s,WHITE); header(s,"05  WHY NOW","레버리지의 '잠식'을 우리는 '수확'으로")
    txt(s,Inches(0.75),Inches(1.05),Inches(9.5),Inches(0.4),[[("같은 변동성, 정반대 매매 — 레버리지 ETF가 '잃는' 그 변동성을 본 펀드는 '수익'으로 전환합니다.",13,DGRAY,True)]])
    pic(s,os.path.join(IMG,"decay_compare.png"),Inches(0.55),Inches(1.55),w=Inches(4.8))
    LRED=RGBColor(0xFC,0xEB,0xE7)
    # 레버리지 박스(잠식=손실)
    rect(s,Inches(5.6),Inches(1.55),Inches(4.45),Inches(1.4),LRED)
    rect(s,Inches(5.6),Inches(1.55),Inches(4.45),Inches(0.5),RED)
    txt(s,Inches(5.75),Inches(1.58),Inches(4.2),Inches(0.45),[[("레버리지 ETF — 변동성에 '잠식'",13,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(5.78),Inches(2.15),Inches(4.1),Inches(0.75),[
        [("매일 고가매수·저가매도(추세 순응 리밸런싱)",11.5,DGRAY,True)],
        [("→ 변동성이 클수록 자산이 녹는다 (기초 제자리인데 −20%)",11,DGRAY,False)]],sp=1.15)
    txt(s,Inches(5.6),Inches(3.02),Inches(4.45),Inches(0.32),[[("▼ 정반대 매매",12,ORANGE,True)]],align=PP_ALIGN.CENTER)
    # 본 펀드 박스(수취=수익)
    rect(s,Inches(5.6),Inches(3.4),Inches(4.45),Inches(1.4),LIGHT)
    rect(s,Inches(5.6),Inches(3.4),Inches(4.45),Inches(0.5),NAVY)
    txt(s,Inches(5.75),Inches(3.43),Inches(4.2),Inches(0.45),[[("변동성 하베스트 펀드 — 변동성을 '수확'",13,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    txt(s,Inches(5.78),Inches(4.0),Inches(4.1),Inches(0.75),[
        [("주가↓ 비중↑(저가매수)·주가↑ 비중↓(고가매도)",11.5,NAVY,True)],
        [("→ 확대된 변동성을 매매 수익으로 전환",11,DGRAY,False)]],sp=1.15)
    rect(s,Inches(0.75),Inches(5.55),Inches(9.3),Inches(1.05),RGBColor(0xFD,0xF1,0xE6))
    txt(s,Inches(1.0),Inches(5.55),Inches(8.9),Inches(1.05),[
        [("핵심: 레버리지가 '잃는' 바로 그 변동성을, 본 펀드는 '수확'한다.",14,ORANGE,True)],
        [("변동성이 클수록(=잠식이 클수록) 본 펀드의 수익 기회는 오히려 커진다.",12,DGRAY,True)]],sp=1.25,anchor=MSO_ANCHOR.MIDDLE)
    footer(s,6)

    # ---- 왜 결론 ----
    s=S(); bg(s,WHITE); header(s,"06  WHY NOW","결론 — 변동성↑ × 방향성↓ 국면의 최적해")
    txt(s,Inches(0.75),Inches(1.12),Inches(9.5),Inches(0.4),[[("변동성은 높고 방향성은 제한적인 지금, 변동성 하베스트 펀드가 가장 잘 맞는 국면입니다.",13,DGRAY,True)]])
    ox,oy,cw,ch=2.0,1.8,3.45,2.05
    quad=[(0,0,ORANGE,"변동성 하베스트 펀드 ★","← 현재 국면",True),
          (1,0,LIGHT,"레버리지·추세추종","방향 확신 필요",False),
          (0,1,LIGHT,"현금성·채권","관망",False),
          (1,1,LIGHT,"인덱스·방향투자","방향 베팅",False)]
    for cx,cy,col,t1,t2,hl in quad:
        rx=Inches(ox+cx*(cw+0.06)); ry=Inches(oy+cy*(ch+0.06))
        rect(s,rx,ry,Inches(cw),Inches(ch),col)
        tc=WHITE if hl else DGRAY
        txt(s,rx+Inches(0.12),ry+Inches(0.12),Inches(cw-0.24),Inches(ch-0.24),[
            [(t1,15 if hl else 13,tc,True)],[(t2,11,(WHITE if hl else GRAY),False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=1.25)
    txt(s,Inches(0.78),Inches(oy),Inches(1.15),Inches(0.35),[[("변동성 ↑",11,NAVY,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.78),Inches(oy+ch+0.06+ch-0.35),Inches(1.15),Inches(0.35),[[("변동성 ↓",11,NAVY,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(ox),Inches(oy+2*ch+0.14),Inches(2*cw+0.06),Inches(0.32),[[("← 방향성 약함           방향성 강함 →",11,NAVY,True)]],align=PP_ALIGN.CENTER)
    rect(s,Inches(0.75),Inches(6.35),Inches(9.3),Inches(0.66),RGBColor(0xFD,0xF1,0xE6))
    txt(s,Inches(0.95),Inches(6.35),Inches(8.9),Inches(0.66),[[("→ 변동성은 '수확'하고 방향성 부재 리스크는 '회피' = 변동성 하베스트 펀드",13,ORANGE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s,7)

    # ---- 전략 ----
    s=S(); bg(s,WHITE); header(s,"07  STRATEGY","전략 — 변동성 매매 + 방향 투자")
    txt(s,Inches(0.75),Inches(1.45),Inches(9.4),Inches(0.4),[[("변동성이 높아 매매 기회는 많고, 방향은 위(상승 참여)라는 두 관점을 하나의 포트폴리오에 결합.",13,DGRAY,True)]])
    boxes=[("① 변동성 매매","저가매수·고가매도를 반복하여 높아진 변동성을 수익으로 전환.",NAVY),
           ("② 방향 투자","상승 구간 추가 수익. 실적·AI·레버리지 수급의 상방에 베팅.",ORANGE),
           ("③ 리스크 관리","최대 편입비 180% 제한 · 급락 시 편입비 축소로 하방 완충.",CYAN)]
    for i,(t,d,c) in enumerate(boxes):
        by=Inches(2.15+i*1.35)
        rect(s,Inches(0.75),by,Inches(0.9),Inches(1.15),c)
        txt(s,Inches(0.75),by,Inches(0.9),Inches(1.15),[[(t.split()[0],20,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        rect(s,Inches(1.75),by,Inches(8.3),Inches(1.15),LIGHT)
        txt(s,Inches(1.95),by+Inches(0.12),Inches(7.9),Inches(0.95),[
            [(t[2:],14,c,True)],[(d,12.5,DGRAY,False)]],sp=1.15,anchor=MSO_ANCHOR.MIDDLE)
    footer(s,8)

    # ---- 운용 방식: 저가매수·고가매도 ----
    s=S(); bg(s,WHITE); header(s,"08  HOW","어떻게 운용하는가 — 저가매수·고가매도")
    txt(s,Inches(0.75),Inches(1.05),Inches(9.5),Inches(0.4),[[("시장이 빠지면 편입을 늘리고(저가매수), 오르면 줄인다(고가매도). 등락(변동성)을 매매로 수취합니다.",13,DGRAY,True)]])
    pic(s,os.path.join(IMG,"oper_flow.png"),Inches(0.6),Inches(1.5),w=Inches(9.3))
    rect(s,Inches(0.75),Inches(6.15),Inches(9.3),Inches(0.6),RGBColor(0xFD,0xF1,0xE6))
    txt(s,Inches(0.95),Inches(6.15),Inches(8.9),Inches(0.6),[[("→ 레버리지 ETF의 '고가매수·저가매도'와 정반대 — 변동성을 잃지 않고 '수취'하는 구조",13,ORANGE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s,9)

    # ---- 수익 구조 ① 성장변동성펀드 (배리어 미터치/터치) ----
    s=S(); bg(s,WHITE); header(s,"09  PAYOFF","수익 구조 — 성장변동성펀드 (변동성 매매 + 상승 참여)")
    txt(s,Inches(0.75),Inches(1.05),Inches(9.5),Inches(0.4),[[("운용 중 -40% 미도달 시 상승에 참여, -40% 도달 시 하방에 노출됩니다. 시뮬레이션 결과(각 점=1회).",12.5,DGRAY,True)]])
    pic(s,os.path.join(IMG,f"scat_{tag}_growth_nt.png"),Inches(0.35),Inches(1.55),w=Inches(4.95))
    pic(s,os.path.join(IMG,f"scat_{tag}_growth_t.png"), Inches(5.45),Inches(1.55),w=Inches(4.95))
    txt(s,Inches(0.75),Inches(6.55),Inches(9.4),Inches(0.4),[
        [(f"σ {int(sig*100)}% 가정 · 점=시뮬 1회 · ",9,GRAY,False),("노란선=구간 평균",9,ORANGE,True),
         (" · 데이터: 자체 시뮬레이션(거래비용·운용오차 포함)",9,GRAY,False)]])
    footer(s,10)

    # ---- 수익 구조 ② 안정변동성펀드 (안정 변동성 매매) ----
    s=S(); bg(s,WHITE); header(s,"10  PAYOFF","수익 구조 — 안정변동성펀드 (안정 변동성 매매)")
    txt(s,Inches(0.75),Inches(1.1),Inches(9.5),Inches(0.4),[[("보합·상승 구간에서 변동성 매매 수익을 안정적으로 확보(상단 고정), 급락 시 손실이 누적됩니다.",13,DGRAY,True)]])
    pic(s,os.path.join(IMG,f"scat_{tag}_stable.png"),Inches(1.55),Inches(1.6),w=Inches(7.7))
    txt(s,Inches(0.75),Inches(6.55),Inches(9.4),Inches(0.4),[
        [(f"σ {int(sig*100)}% 가정 · 점=시뮬 1회 · ",9,GRAY,False),("노란선=구간 평균",9,ORANGE,True),
         (" · 데이터: 자체 시뮬레이션(거래비용·운용오차 포함)",9,GRAY,False)]])
    footer(s,11)

    # ---- 리스크 ----
    s=S(); bg(s,WHITE); header(s,"11  RISK","리스크 및 유의사항")
    bullets(s,Inches(0.75),Inches(1.6),Inches(9.3),Inches(5),[
        ("■","하방 손실 위험 (최대 위험요소): 전략 특성상 기초자산 급락 시 손실이 누적되며 원금의 상당 부분 손실이 가능합니다.",NAVY,True),
        ("■","운용오차: 이산 리밸런싱·급변동·갭 하락 시 이론 수익과 실현손익이 크게 벌어질 수 있습니다(시뮬레이션상 이탈 사례 확인).",DGRAY,False),
        ("■","편입비 상한(180%) 효과: 급변동 구간에서 편입비 조절이 제한되어 추가 오차가 발생할 수 있습니다.",DGRAY,False),
        ("■","급락 국면 위험: 급격한 하락 시 하방 완충이 제한되어 손실에 노출될 수 있습니다.",DGRAY,False),
        ("■",("단일종목 집중 위험(AI Top2)" if not is_index else "지수 변동 위험"),DGRAY,True) if not is_index else ("■","지수 변동·시장 위험: 시장 전반의 급락 국면에서 손실이 발생할 수 있습니다.",DGRAY,False),
        ("■","변동성 축소 위험: 시장 변동성이 가정보다 낮아지면 변동성 매매 수익 기회가 줄어 기대수익에 미달할 수 있습니다.",DGRAY,False),
    ],sz=12.5,gap=1.3)
    if not is_index:
        txt(s,Inches(1.05),Inches(4.15),Inches(9),Inches(0.4),[[("삼성전자·SK하이닉스 2종목에 집중되어 개별기업·반도체 사이클 리스크에 크게 노출됩니다.",11.5,GRAY,False)]])
    rect(s,Inches(0.75),Inches(5.55),Inches(9.3),Inches(1.05),LIGHT)
    txt(s,Inches(1.0),Inches(5.68),Inches(8.9),Inches(0.85),[
        [("본 상품은 1등급(매우 높은 위험) 상품으로 원금의 전부(0~100%) 손실이 발생할 수 있으며, 그 손실은 투자자에게 귀속됩니다.",11.5,RED,True)],
        [("투자 전 상품설명서·약관·계약권유문서를 반드시 확인하시기 바랍니다.",10.5,GRAY,False)]],sp=1.2)
    footer(s,12)

    # ---- 운용부서/컴플 ----
    s=S(); bg(s,NAVY)
    rect(s,Inches(0.9),Inches(1.4),Inches(0.14),Inches(0.55),ORANGE)
    txt(s,Inches(1.15),Inches(1.4),Inches(9),Inches(0.6),[[("운용부서 및 유의사항",22,WHITE,True)]])
    txt(s,Inches(1.15),Inches(2.6),Inches(9),Inches(2),[
        [("AI금융공학운용부문 전략운용본부",15,CYAN,True)],
        [("Mission: 투자자 입장에서 장기 플러스 수익 달성",13,WHITE,False)],
        [("정량적 분석과 정성적 리서치의 균형 · One Team Approach",12,RGBColor(0xC9,0xD6,0xE5),False)],
        [("커버드콜·국내인덱스+·절대수익·구조화·Custom Mandate 등 운용",12,RGBColor(0xC9,0xD6,0xE5),False)]],sp=1.4)
    txt(s,Inches(1.15),Inches(5.3),Inches(8.9),Inches(1.6),[
        [("· 본 자료는 컴플라이언스 검토 전 내부 검토용 초안이며 대외 배포를 금합니다.",10,LORANGE,False)],
        [("· 랩/펀드 계약은 예금자보호법에 따라 보호되지 않으며, 자산가격·환율·신용 변동에 따라 원금손실(0~100%)이 발생할 수 있습니다.",10,RGBColor(0xC9,0xD6,0xE5),False)],
        [("· 상기 시뮬레이션은 가정에 의거한 예시로 수익을 보장하지 않으며 실제 결과와 다를 수 있습니다.",10,RGBColor(0xC9,0xD6,0xE5),False)]],sp=1.3)
    footer(s,13,"")

    prs.save(path)
    print("saved:",path)

base=os.path.join(os.path.dirname(os.path.abspath(__file__)),"")
build(base+"변동성하베스트펀드_KOSPI200_초안.pptx","KOSPI200","코스피200 지수",0.50,True,  os.path.join(IMG,"vol_kospi200.png"),"kospi")
build(base+"변동성하베스트펀드_AITop2_초안.pptx","AI Top2 (삼성전자·SK하이닉스)","AI Top2 — 삼성전자·SK하이닉스 (동일가중)",0.60,False, os.path.join(IMG,"vol_semi.png"),"top2")
print("done")
