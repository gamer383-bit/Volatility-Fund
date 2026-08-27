# -*- coding: utf-8 -*-
"""P3 편입비(100~200%) 표기 추가 · P4 국면별 대비 페이지 (KOSPI 실측, data 폴더 DataGuide)
- 추세장(2025.09.26~2026.01.30, KOSPI +54%): 기존 레버리지 유리
- 출렁 상승장(2026.04.17~2026.08.21, KOSPI +12%·변동성 70%): 스마트 레버리지 유리
- 기존 = 일간 1.5배 매일 재조정 · 스마트 = 편입비 100~200% 범위 내 방치(범위 이탈 시 경계로 복원) · 비용 미반영 개념 예시
"""
import sys, io, os, platform
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
import numpy as np, pandas as pd
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
plt.rcParams['font.family']=('Malgun Gothic' if platform.system()=='Windows' else 'AppleGothic')
plt.rcParams['axes.unicode_minus']=False
NAVY='#043B72'; ORANGE='#F58220'; GRAY='#84888B'
os.chdir(os.path.join("C:/Users/gamer38/Documents/Claude/Projects","변동성펀드","제안서"))
f=os.path.join("..","data","삼성전자_하이닉스_코스피_코스피200_코스닥150_최근5년_주가데이터.xlsx")
raw=pd.read_excel(f,sheet_name=1,header=None,skiprows=14)
ks=pd.Series(raw[1].astype(float).values,index=pd.to_datetime(raw[0])).sort_index().dropna()
rday=0.025/252

def sim_paths(seg):
    rr=seg.pct_change().fillna(0.0)
    vb=[1.0]; vs=[1.0]; E=1.5
    for x in rr.values[1:]:
        vb.append(vb[-1]*(1+1.5*x-0.5*rday))
        L=E/vs[-1]
        v=vs[-1]*(1+L*x+(1-L)*rday)
        E*=1+x
        L=E/v
        if L>2.0: E=2.0*v
        elif L<1.0: E=1.0*v
        vs.append(v)
    return np.array(vb),np.array(vs)

def style(ax):
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.grid(axis='y',color='#e6ebf1',lw=0.7); ax.set_axisbelow(True)
    ax.tick_params(labelsize=8.5)

wins=[("추세가 강하고 지속되는 장 : 기존 레버리지 유리","2025-09-26","2026-01-30"),
      ("변동성이 크고 출렁이며 오르는 장 : 스마트 유리","2026-04-17","2026-08-21")]
fig,axes=plt.subplots(1,2,figsize=(10.4,3.4),dpi=150)
for ax,(tt,a,b) in zip(axes,wins):
    seg=ks[(ks.index>=pd.Timestamp(a))&(ks.index<=pd.Timestamp(b))]
    vb,vs=sim_paths(seg)
    ix=seg/seg.iloc[0]*100
    ax.plot(seg.index,ix.values,color=GRAY,lw=1.4,label='KOSPI')
    ax.plot(seg.index,vb*100,color='#C0392B',lw=1.7,label='기존 레버리지 1.5 (매일 재조정)')
    ax.plot(seg.index,vs*100,color=ORANGE,lw=2.0,label='스마트 레버리지 1.5')
    span=max(vb.max(),vs.max(),ix.max()/100)*100-min(vb.min(),vs.min(),ix.min()/100)*100
    labs=sorted([(vb[-1]*100,'#C0392B'),(vs[-1]*100,ORANGE),(float(ix.iloc[-1]),GRAY)],key=lambda t:-t[0])
    pos=[]
    for v,_ in labs: pos.append(v if not pos else min(v,pos[-1]-span*0.09))
    for (v,c),py in zip(labs,pos):
        ax.annotate(f"{v-100:+.1f}%",xy=(seg.index[-1],py),xytext=(5,-3),
                    textcoords='offset points',fontsize=9.5,fontweight='bold',color=c)
    win=vs[-1]-vb[-1]
    lab=f"스마트 {win*100:+.1f}%p" if win>0 else f"기존 {-win*100:+.1f}%p"
    ax.annotate(lab,xy=(0.04,0.90),xycoords='axes fraction',fontsize=10.5,fontweight='bold',
                color=(ORANGE if win>0 else '#C0392B'))
    g=float(seg.iloc[-1]/seg.iloc[0]-1); vol=float(seg.pct_change().std())*np.sqrt(252)
    ax.set_title(f"{tt}\n({a[:7].replace('-','.')} ~ {b[:7].replace('-','.')} · KOSPI {g*100:+.0f}% · 변동성 {vol*100:.0f}%)",
                 fontsize=10.5,fontweight='bold',color='#222')
    ax.legend(fontsize=7.5,loc='upper left',frameon=False,bbox_to_anchor=(0.0,0.88))
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%y.%m'))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=1)); style(ax)
    ax.set_xlim(seg.index[0],seg.index[-1]+pd.Timedelta(days=16))
    print(f"[{tt}] KOSPI {g*100:+.1f}% vol {vol*100:.0f}% | 기존 {vb[-1]*100-100:+.1f}% vs 스마트 {vs[-1]*100-100:+.1f}%")
plt.tight_layout()
plt.savefig(os.path.join('img','lev15_regime.png')); plt.close()
print("차트 저장: lev15_regime.png")

# ---- 덱 반영 ----
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
DST="Tech_TOP10_스마트액티브레버리지15_v1_20260812.pptx"
prs=Presentation(DST)
DARK=RGBColor(0x22,0x22,0x22); NAVY_=RGBColor(0x04,0x3B,0x72); GRAY_=RGBColor(0x84,0x88,0x8B)
ORANGE_=RGBColor(0xE8,0x72,0x0C)
FT_B="KoPub돋움체_Pro Bold"; FT_M="KoPub돋움체_Pro Medium"
SW=prs.slide_width/914400; SH=prs.slide_height/914400
def tbox(s,L,T,W,H,text,size,bold,color,font=None,sp=1.15):
    b=s.shapes.add_textbox(Inches(L),Inches(T),Inches(W),Inches(H))
    tf=b.text_frame; tf.word_wrap=True
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
        if font: r.font.name=font
        p.line_spacing=sp
    return b
def para_replace(p,a,b):
    runs=p.runs; full=''.join(r.text for r in runs)
    if a not in full: return 0
    st=full.index(a); en=st+len(a); pos=0; first=True
    for r in runs:
        ln=len(r.text); s0,e0=pos,pos+ln
        if e0<=st or s0>=en: pos=e0; continue
        pre=r.text[:max(st-s0,0)]; post=r.text[max(en-s0,0):] if en-s0<ln else ''
        if first: r.text=pre+b+post; first=False
        else: r.text=pre+post
        pos=e0
    return 1

# P3: 편입비 범위 표기
s3=prs.slides[2]; n=0
OLD="기준 레버리지는 1.5배. 노출이 목표 범위 안에 있으면 매매하지 않습니다 — 레버리지의 기대수익은 그대로 가져갑니다."
NEW="기준 레버리지는 1.5배이며, 편입비는 최소 100% ~ 최대 200% 범위에서 유지합니다. 노출이 범위 안에 있으면 매매하지 않아 레버리지의 기대수익은 그대로 가져갑니다."
for sh in s3.shapes:
    if getattr(sh,'has_text_frame',False) and sh.has_text_frame:
        for p in sh.text_frame.paragraphs: n+=para_replace(p,OLD,NEW)
print(f"P3 편입비 표기 {n}건")

# P4: 국면별 대비 페이지 재작성
s4=prs.slides[3]
for sh in list(s4.shapes): s4.shapes._spTree.remove(sh._element)
tbox(s4,0.54,0.42,SW-1.1,0.45,"1. 스마트 레버리지란 : 국면별 — 기존 레버리지와 무엇이 다른가",21,True,DARK,FT_B)
tbox(s4,0.55,0.98,SW-1.1,0.6,
     "한 방향으로 강하게 지속되는 추세장에서는 매일 배수를 유지하는 기존 레버리지가 유리합니다. "
     "반면 변동성이 크고 출렁이며 오르는 장에서는 고가 매도·저가 매수가 작동하는 스마트 레버리지가 유리합니다 — KOSPI 실측으로 비교했습니다.",13.5,None,DARK)
s4.shapes.add_picture(os.path.join('img','lev15_regime.png'),Inches(0.45),Inches(2.05),width=Inches(9.95))
tbox(s4,0.8,5.75,SW-1.7,0.5,"어느 한쪽이 항상 이기는 방식은 없습니다 — 최근처럼 변동성이 큰 시장일수록 스마트 방식의 강점이 커집니다.",14,True,ORANGE_,FT_B)
tbox(s4,0.55,SH-0.62,SW-1.1,0.4,
     "※ KOSPI 종가지수 실측(DataGuide, ~2026.08.26) 기반 개념 시뮬레이션 · 기존 = 일간 1.5배 매일 재조정 · 스마트 = 편입비 100~200% 범위 내 운용 가정 · 거래비용 미반영 · 상기 자료는 이해를 돕기 위한 예시이며 실제 성과와 다를 수 있습니다.",8,False,GRAY_,FT_M)
prs.save(DST)
print("P3·P4 저장 완료")
