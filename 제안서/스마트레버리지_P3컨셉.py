# -*- coding: utf-8 -*-
"""P3 조정: ①평균 1.5배 ②상승 시 매도·하락 시 매수(잠식 제거) + 기존 레버리지 대비 개념도"""
import sys, io, os
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
import numpy as np
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import platform
plt.rcParams['font.family']=('Malgun Gothic' if platform.system()=='Windows' else 'AppleGothic')
plt.rcParams['axes.unicode_minus']=False
NAVY='#043B72'; RED='#C0392B'; BLUE='#1F6FB2'; GRAY='#84888B'
os.chdir(os.path.join("C:/Users/gamer38/Documents/Claude/Projects","변동성펀드","제안서"))

# ---- 개념도: 같은 파동 위에서 매매 방향만 반대 ----
x=np.linspace(0,4*np.pi,400)
y=np.sin(x-np.pi/2)          # 저점에서 시작해 고점-저점-고점
peaks=[np.pi/2+np.pi*0+np.pi/2, ]  # compute directly
px_=[np.pi, 3*np.pi]          # y=sin(x-pi/2) 최고점: x-pi/2=pi/2 → x=pi, x=3pi
tx_=[2*np.pi, 4*np.pi]        # 최저점: x-pi/2=3pi/2 → x=2pi, 4pi
def wave(ax):
    ax.plot(x,y,color=GRAY,lw=2.2,zorder=1)
    ax.set_xlim(-0.4,4*np.pi+0.4); ax.set_ylim(-2.55,2.1)
    ax.axis('off')
fig,(a1,a2)=plt.subplots(1,2,figsize=(10.2,2.9),dpi=150)
# 좌: 기존 레버리지 — 고점 매수(▲빨강), 저점 매도(▼파랑)
wave(a1)
a1.scatter(px_, np.sin(np.array(px_)-np.pi/2)+0.13, marker='^', s=150, color=RED, zorder=3)
a1.scatter(tx_[:1]+[0.0], [np.sin(2*np.pi-np.pi/2)-0.13, np.sin(0-np.pi/2)-0.13], marker='v', s=150, color=BLUE, zorder=3)
for xx in px_: a1.annotate('매수',xy=(xx,np.sin(xx-np.pi/2)+0.35),ha='center',fontsize=11,fontweight='bold',color=RED)
for xx in [0.0,2*np.pi]: a1.annotate('매도',xy=(xx,np.sin(xx-np.pi/2)-0.62),ha='center',fontsize=11,fontweight='bold',color=BLUE)
a1.set_title("기존 레버리지 ETF : 오르면 매수, 내리면 매도",fontsize=12.5,fontweight='bold',color='#222',pad=10)
a1.annotate("고가에 사고 저가에 판다 → 왕복할 때마다 손실 (변동성 잠식)",xy=(2*np.pi,-2.35),ha='center',fontsize=10.5,fontweight='bold',color=RED)
# 우: 스마트 레버리지 — 고점 매도(▼파랑), 저점 매수(▲빨강)
wave(a2)
a2.scatter(px_, np.sin(np.array(px_)-np.pi/2)+0.13, marker='v', s=150, color=BLUE, zorder=3)
a2.scatter([0.0,2*np.pi], [np.sin(0-np.pi/2)-0.13,np.sin(2*np.pi-np.pi/2)-0.13], marker='^', s=150, color=RED, zorder=3)
for xx in px_: a2.annotate('매도',xy=(xx,np.sin(xx-np.pi/2)+0.35),ha='center',fontsize=11,fontweight='bold',color=BLUE)
for xx in [0.0,2*np.pi]: a2.annotate('매수',xy=(xx,np.sin(xx-np.pi/2)-0.62),ha='center',fontsize=11,fontweight='bold',color=RED)
a2.set_title("스마트 레버리지 1.5 : 오르면 매도, 내리면 매수",fontsize=12.5,fontweight='bold',color=NAVY,pad=10)
a2.annotate("고가에 팔고 저가에 산다 → 변동성 잠식을 없앤다",xy=(2*np.pi,-2.35),ha='center',fontsize=10.5,fontweight='bold',color=NAVY)
plt.tight_layout()
plt.savefig(os.path.join('img','lev15_concept.png')); plt.close()
print("개념도 저장: lev15_concept.png")

# ---- P3 재작성 ----
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
DST="Tech_TOP10_스마트액티브레버리지15_v1_20260812.pptx"
prs=Presentation(DST)
DARK=RGBColor(0x22,0x22,0x22); NAVY_=RGBColor(0x04,0x3B,0x72); GRAY_=RGBColor(0x84,0x88,0x8B)
FT_B="KoPub돋움체_Pro Bold"; FT_M="KoPub돋움체_Pro Medium"
SW=prs.slide_width/914400; SH=prs.slide_height/914400
def tbox(s,L,T,W,H,text,size,bold,color,font=None,sp=1.15):
    b=s.shapes.add_textbox(Inches(L),Inches(T),Inches(W),Inches(H))
    tf=b.text_frame; tf.word_wrap=True
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
        if font: r.font.name=font
        p.line_spacing=sp
    return b
s=prs.slides[2]
for sh in list(s.shapes): s.shapes._spTree.remove(sh._element)
tbox(s,0.54,0.42,SW-1.1,0.45,"1. 스마트 레버리지란 : '매일'이 아니라 '평균'으로 1.5배",21,True,DARK,FT_B)
tbox(s,0.55,0.98,SW-1.1,0.6,"'평균 1.5배' 노출을 유지하면서, 오르면 팔고 내리면 사는 운용으로 기존 레버리지의 변동성 잠식을 없애는 방식입니다.",13.5,None,DARK)
tbox(s,0.8,1.85,SW-1.7,0.35,"① 평균 1.5배 노출",15,True,NAVY_,FT_B)
tbox(s,0.8,2.27,SW-1.7,0.9,"기준 레버리지는 1.5배. 노출이 목표 범위 안에 있으면 매매하지 않습니다 — 레버리지의 기대수익은 그대로 가져갑니다.",12,None,DARK,sp=1.25)
tbox(s,0.8,2.95,SW-1.7,0.35,"② 오르면 매도, 내리면 매수",15,True,NAVY_,FT_B)
tbox(s,0.8,3.37,SW-1.7,0.9,"기존 레버리지 ETF는 오르면 사고 내리면 파는 일일 재조정으로 왕복 구간마다 손실이 쌓입니다. "
     "스마트 레버리지는 반대로 상승 시 매도(이익 실현), 하락 시 매수(저가 매수)를 하여 변동성 잠식 효과를 없앱니다.",12,None,DARK,sp=1.25)
s.shapes.add_picture(os.path.join('img','lev15_concept.png'),Inches(0.65),Inches(4.35),width=Inches(9.55))
tbox(s,0.55,SH-0.62,SW-1.1,0.4,"※ 상기 자료는 이해를 돕기 위한 개념 예시이며, 시뮬레이션·운용은 가정과 시장 상황에 따라 달라질 수 있습니다. 레버리지 상품은 원금 손실 위험이 큽니다.",8,False,GRAY_,FT_M)
prs.save(DST)
print("P3 저장 완료")
