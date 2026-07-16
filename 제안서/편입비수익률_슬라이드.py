# -*- coding: utf-8 -*-
"""AITop2 v1: 10·11페이지 재구성
- 기존 편입비 슬라이드 2장(10·11) 삭제 → 10=안정형, 11=성장형
- 각 페이지: 좌 편입비 테이블 + 우 예상수익률 테이블 (지수 60~140 × 잔존 12/9/6/3개월)
- 60행: 성장형 낙아웃 후 실제 편입비(≈100%) 표기, 헤더 '지수 / 잔존만기'
- 예상수익률 ≥15% 셀 음영, 변동성 변화 유의문구
"""
import sys, io, os, math
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from scipy.special import ndtr as N
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

R,Q,SIG=0.025,0.025,0.60   # 금리=배당=2.5% (사용자 지정)
KPUT,KKO,H,K1,K2=100.,100.,60.,110.,150.
CAP_G, CAP_S = 1.8, 1.0     # 성장형 180% / 안정형 100%
NAVY=RGBColor(0x04,0x3B,0x72); ORANGE=RGBColor(0xF5,0x82,0x20); GRAY=RGBColor(0x84,0x88,0x8B)
DGRAY=RGBColor(0x3A,0x3A,0x3A); LIGHT=RGBColor(0xEE,0xF2,0xF7); WHITE=RGBColor(0xFF,0xFF,0xFF)
HL=RGBColor(0xFB,0xDD,0xB9)   # ≥15% 음영 (연주황)
FT="맑은 고딕"

def bs_delta(typ,S,K,T):
    sq=SIG*math.sqrt(T); d1=(math.log(S/K)+(R-Q+0.5*SIG*SIG)*T)/sq
    eqt=math.exp(-Q*T)
    return eqt*N(d1) if typ=='c' else eqt*(N(d1)-1)
def bs_price(typ,S,K,T):
    sq=SIG*math.sqrt(T); d1=(math.log(S/K)+(R-Q+0.5*SIG*SIG)*T)/sq; d2=d1-sq
    eqt=math.exp(-Q*T); ert=math.exp(-R*T)
    if typ=='c': return S*eqt*N(d1)-K*ert*N(d2)
    return K*ert*N(-d2)-S*eqt*N(-d1)
def din_put(S,K,Hb,T):
    b=R-Q; sq=SIG*math.sqrt(T); mu=(b-0.5*SIG*SIG)/(SIG*SIG); phi=-1.; eta=1.
    x2=math.log(S/Hb)/sq+(1+mu)*sq; y1=math.log(Hb*Hb/(S*K))/sq+(1+mu)*sq; y2=math.log(Hb/S)/sq+(1+mu)*sq
    ebr=math.exp((b-R)*T); ert=math.exp(-R*T); pw1=(Hb/S)**(2*(mu+1)); pw2=(Hb/S)**(2*mu)
    B=phi*S*ebr*N(phi*x2)-phi*K*ert*N(phi*x2-phi*sq)
    C=phi*S*ebr*pw1*N(eta*y1)-phi*K*ert*pw2*N(eta*y1-eta*sq)
    D=phi*S*ebr*pw1*N(eta*y2)-phi*K*ert*pw2*N(eta*y2-eta*sq)
    return B-C+D
def dop(S,K,Hb,T):
    if S<=Hb: return 0.0
    return max(bs_price('p',S,K,T)-din_put(S,K,Hb,T),0.0)
def ko_delta(S,K,Hb,T):
    if S<=Hb: return 0.0
    h=max(S*1e-4,1e-6)
    return (dop(S+h,K,Hb,T)-dop(max(S-h,Hb+1e-9),K,Hb,T))/(2*h)
def cs_price(S,T): return bs_price('c',S,K1,T)-bs_price('c',S,K2,T)
def cs_delta(S,T): return bs_delta('c',S,K1,T)-bs_delta('c',S,K2,T)

P0=bs_price('p',100.,KPUT,1.0)
V0_S=-P0
V0_G=-P0+dop(100.,KKO,H,1.0)+cs_price(100.,1.0)

def w_stable(S,tau): return min(max(-bs_delta('p',S,KPUT,tau),0.0),CAP_S)*100
def w_growth(S,tau):
    if S<=H:  # 낙아웃: KO 소멸, 풋매도+콜스프레드 델타 (사실상 주식형 ≈100%)
        d=-bs_delta('p',S,KPUT,tau)+cs_delta(S,tau)
    else:
        d=-bs_delta('p',S,KPUT,tau)+ko_delta(S,KKO,H,tau)+cs_delta(S,tau)
    return min(max(d,0.0),CAP_G)*100
def ret_stable(S,tau):
    return (V0_S*math.exp(R*(1.0-tau)) - (-bs_price('p',S,KPUT,tau)))*-1/100*100  # = (P0 e^{rΔt} − Put)/100
def ret_stable2(S,tau):  # 명시적: (−Put(m,τ)) − V0·e^{rΔt}
    return ((-bs_price('p',S,KPUT,tau)) - V0_S*math.exp(R*(1.0-tau)))
def ret_growth(S,tau):
    if S<=H: V=-bs_price('p',S,KPUT,tau)+cs_price(S,tau)
    else:    V=-bs_price('p',S,KPUT,tau)+dop(S,KKO,H,tau)+cs_price(S,tau)
    return V - V0_G*math.exp(R*(1.0-tau))

LEVELS=[140,130,120,110,100,90,80,70,60]
TAUS=[(12,1.0),(9,0.75),(6,0.5),(3,0.25)]
HDR=["지수 / 잔존만기"]+[f"{m}개월" for m,_ in TAUS]

def build(fn_w,fn_r):
    wrows=[HDR[:]]; rrows=[HDR[:]]; shade=[[False]*5]
    for lv in LEVELS:
        wr=[f"{lv}"]; rr=[f"{lv}"]; sh=[False]
        for _,tau in TAUS:
            wr.append(f"{fn_w(float(lv),tau):.0f}%")
            rv=fn_r(float(lv),tau)
            rr.append(f"{rv:+.1f}%"); sh.append(rv>=15.0)
        wrows.append(wr); rrows.append(rr); shade.append(sh)
    return wrows,rrows,shade

sw,sr,ss=build(w_stable,ret_stable2)
gw,gr,gs=build(w_growth,ret_growth)
print("=== 안정형 편입비 ==="); [print(" | ".join(r)) for r in sw]
print("=== 안정형 예상수익률 ==="); [print(" | ".join(r)) for r in sr]
print("=== 성장형 편입비 ==="); [print(" | ".join(r)) for r in gw]
print("=== 성장형 예상수익률 ==="); [print(" | ".join(r)) for r in gr]

PATH="C:/Users/gamer38/Documents/Claude/Projects/변동성펀드/제안서/변동성하베스트펀드_AITop2_v1_20260715.pptx"
prs=Presentation(PATH)

# ---- 삭제 대상(기존 편입비 슬라이드, 현재 10·11페이지 = idx 9,10) 핸들 확보 (삭제는 '추가 후'에!) ----
# 새 슬라이드를 먼저 추가해야 파트명이 유일하게 배정됨(추가→재배치→삭제 순서, 파트명 충돌 방지)
_old=[list(prs.slides._sldIdLst)[9], list(prs.slides._sldIdLst)[10]]

def make_slide(pos, kicker, title, sub, name, color, wrows, rrows, shade, extra_note):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    def rect(x,y,w,h,c):
        r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
        r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); r.shadow.inherit=False
        return r
    def txt(x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0):
        tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
        tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
        if isinstance(runs[0],tuple): runs=[runs]
        for i,para in enumerate(runs):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.alignment=align; p.line_spacing=sp
            for (t,sz,col,b) in para:
                r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=col; r.font.name=FT
        return tb
    def table(x,y,w,rows,hfill,shades=None,fs=10):
        nr,nc=len(rows),len(rows[0]); rh=Inches(0.36)
        gt=s.shapes.add_table(nr,nc,x,y,w,rh*nr).table
        cw=[Inches(1.35)]+[(w-Inches(1.35))//(nc-1)]*(nc-1)
        for j,c in enumerate(cw): gt.columns[j].width=c
        for i,row in enumerate(rows):
            gt.rows[i].height=rh
            for j,val in enumerate(row):
                c=gt.cell(i,j); c.margin_left=Pt(3); c.margin_right=Pt(3); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
                c.vertical_anchor=MSO_ANCHOR.MIDDLE
                p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=str(val)
                r.font.name=FT; r.font.size=Pt(fs if j>0 else fs-0.5); p.alignment=PP_ALIGN.CENTER
                if i==0:
                    c.fill.solid(); c.fill.fore_color.rgb=hfill; r.font.color.rgb=WHITE; r.font.bold=True
                else:
                    fillc=(LIGHT if i%2==0 else WHITE)
                    if shades and shades[i][j]: fillc=HL
                    c.fill.solid(); c.fill.fore_color.rgb=fillc
                    r.font.color.rgb=DGRAY; r.font.bold=(j==0 or row[0]=="100")
        return gt
    bg=rect(0,0,prs.slide_width,prs.slide_height,WHITE)
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2,bg._element)
    rect(Inches(0.5),Inches(0.55),Inches(0.14),Inches(0.62),ORANGE)
    txt(Inches(0.75),Inches(0.5),Inches(9),Inches(0.35),[[(kicker,12,ORANGE,True)]])
    txt(Inches(0.75),Inches(0.8),Inches(9.5),Inches(0.55),[[(title,23,NAVY,True)]])
    txt(Inches(0.75),Inches(1.42),Inches(9.5),Inches(0.4),[[(sub,12.5,DGRAY,True)]])
    txt(Inches(0.6),Inches(1.95),Inches(4.8),Inches(0.33),[[("실제 편입비",13.5,color,True)]])
    table(Inches(0.6),Inches(2.3),Inches(4.8),wrows,color)
    txt(Inches(5.65),Inches(1.95),Inches(4.8),Inches(0.33),[[("예상수익률",13.5,color,True),("  (원금 대비 · ",10,GRAY,False),("15% 이상 음영",10,ORANGE,True),(")",10,GRAY,False)]])
    table(Inches(5.65),Inches(2.3),Inches(4.75),rrows,color,shades=shade)
    txt(Inches(0.6),Inches(6.15),Inches(9.6),Inches(0.62),[
        [(extra_note,9.5,GRAY,False)],
        [("· 변동성이 변할 경우 예상 편입비와 예상수익률은 달라질 수 있습니다.",10,DGRAY,True)]],sp=1.25)
    txt(Inches(0.5),Inches(7.12),Inches(8.6),Inches(0.32),[[("※ 상기 자료는 이해를 돕기 위한 예시이며 실제 투자내용을 의미하지 않습니다. 시뮬레이션은 가정에 따른 것으로 실제와 다를 수 있습니다.",7,GRAY,False)]])
    ids2=list(prs.slides._sldIdLst)
    prs.slides._sldIdLst.remove(ids2[-1]); prs.slides._sldIdLst.insert(pos,ids2[-1])

make_slide(9,"08  HOW","지수대별 실제 편입비·예상수익률 — 안정변동성펀드",
    "리밸런싱일 지수=100 기준 · 잔존만기별 편입비와 원금 대비 예상수익률 · 최대 편입비 100%",
    "안정",NAVY,sw,sr,ss,
    "· 변동성 60% 가정, 지수 10포인트 단위(60~140), 잔존만기 12·9·6·3개월.")
make_slide(10,"08  HOW","지수대별 실제 편입비·예상수익률 — 성장변동성펀드",
    "리밸런싱일 지수=100 기준 · 잔존만기별 편입비와 원금 대비 예상수익률 · 최대 편입비 180%",
    "성장",ORANGE,gw,gr,gs,
    "· 변동성 60% 가정, 지수 10포인트 단위(60~140). 지수 60 도달 시 하방 완충이 소멸되어 주식형과 유사하게 운용(실제 편입비 약 100%).")

# ---- 이제 기존 편입비 슬라이드 2장 삭제 (새 슬라이드 추가 후라 파트명 충돌 없음) ----
xml=prs.slides._sldIdLst
for sid in _old:
    rId=sid.get(qn('r:id'))
    prs.part.drop_rel(rId); xml.remove(sid)
print("기존 편입비 슬라이드 2장 삭제 (추가 후 삭제 순서)")

try:
    prs.save(PATH); print("[OK] 10(안정형)·11(성장형) 페이지 재구성 완료")
except PermissionError:
    print("[잠김] PowerPoint에 열려 있어 저장 불가")
