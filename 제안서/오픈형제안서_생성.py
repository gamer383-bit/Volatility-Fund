# -*- coding: utf-8 -*-
"""오픈형 변동성 하베스트 레버리지 제안서 (별도 덱) — 방법론 + 실측 백테스트
용어 방침: 대외 순화 톤(변동성 매매·저가매수/고가매도·옵션이론 기반 모형), 내부 검토용 초안"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY=RGBColor(0x04,0x3B,0x72); ORANGE=RGBColor(0xF5,0x82,0x20); CYAN=RGBColor(0x00,0xA9,0xCE)
GRAY=RGBColor(0x84,0x88,0x8B); DGRAY=RGBColor(0x3A,0x3A,0x3A); LIGHT=RGBColor(0xEE,0xF2,0xF7)
WHITE=RGBColor(0xFF,0xFF,0xFF); LORANGE=RGBColor(0xF0,0xB2,0x6B); RED=RGBColor(0xD0,0x2F,0x00)
FT="맑은 고딕"; SW,SH=Inches(10.8),Inches(7.5)
BASE=os.path.dirname(os.path.abspath(__file__))
OIMG=os.path.join(os.path.dirname(BASE),'오픈형','img')

prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
blank=prs.slide_layouts[6]
def S(): return prs.slides.add_slide(blank)
def bg(s,c):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,SW,SH)
    r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
def rect(s,x,y,w,h,c):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); r.shadow.inherit=False
    return r
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=1.0):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=sp
        for (t,sz,col,b) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=col; r.font.name=FT
    return tb
def bullets(s,x,y,w,h,items,sz=12.5,gap=1.28):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(mark,t,c,b) in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=gap; p.space_after=Pt(3)
        r=p.add_run(); r.text=(mark+" " if mark else "")+t
        r.font.size=Pt(sz); r.font.color.rgb=c; r.font.bold=b; r.font.name=FT
    return tb
def footer(s,pg,note="※ 상기 자료는 이해를 돕기 위한 예시이며 실제 투자내용을 의미하지 않습니다. 시뮬레이션·백테스트는 가정에 따른 것으로 실제와 다를 수 있습니다."):
    txt(s,Inches(0.5),Inches(7.12),Inches(8.6),Inches(0.32),[[(note,7,GRAY,False)]])
    txt(s,Inches(9.9),Inches(7.12),Inches(0.6),Inches(0.32),[[(str(pg),9,GRAY,True)]],align=PP_ALIGN.RIGHT)
def header(s,k,t):
    rect(s,Inches(0.5),Inches(0.55),Inches(0.14),Inches(0.62),ORANGE)
    txt(s,Inches(0.75),Inches(0.5),Inches(9),Inches(0.35),[[(k,12,ORANGE,True)]])
    txt(s,Inches(0.75),Inches(0.8),Inches(9.6),Inches(0.55),[[(t,22,NAVY,True)]])
def table(s,x,y,w,rows,colw,hfill=NAVY,fs=11,rh=Inches(0.38),boldrows=()):
    nr,nc=len(rows),len(rows[0])
    gt=s.shapes.add_table(nr,nc,x,y,w,rh*nr).table
    for j,cw in enumerate(colw): gt.columns[j].width=cw
    for i,row in enumerate(rows):
        gt.rows[i].height=rh
        for j,val in enumerate(row):
            c=gt.cell(i,j); c.margin_left=Pt(4); c.margin_right=Pt(4); c.margin_top=Pt(1); c.margin_bottom=Pt(1)
            c.vertical_anchor=MSO_ANCHOR.MIDDLE
            p=c.text_frame.paragraphs[0]; r=p.add_run(); r.text=str(val)
            r.font.name=FT; r.font.size=Pt(fs)
            if i==0:
                c.fill.solid(); c.fill.fore_color.rgb=hfill; r.font.color.rgb=WHITE; r.font.bold=True; p.alignment=PP_ALIGN.CENTER
            else:
                c.fill.solid(); c.fill.fore_color.rgb=(LIGHT if i%2==0 else WHITE)
                r.font.color.rgb=DGRAY; r.font.bold=(j==0 or i in boldrows)
                p.alignment=(PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
def pic(s,path,x,y,w): return s.shapes.add_picture(path,x,y,width=w)

# ---- 1 표지 ----
s=S(); bg(s,NAVY)
rect(s,0,Inches(2.55),SW,Inches(0.06),ORANGE)
txt(s,Inches(0.9),Inches(1.4),Inches(9),Inches(0.5),[[("미래에셋자산운용  AI금융공학운용부문",14,CYAN,True)]])
txt(s,Inches(0.9),Inches(2.75),Inches(9.6),Inches(1.5),[
    [("오픈형 변동성 하베스트 레버리지",38,WHITE,True)],
    [("KOSPI200 · 삼성전자 · SK하이닉스 시리즈",22,LORANGE,True)]],sp=1.1)
txt(s,Inches(0.9),Inches(5.1),Inches(9.2),Inches(0.7),[
    [("BM 150% · 편입비 100~200% · 상승 시 줄이고 하락 시 늘리는 변동성 매매로",15,RGBColor(0xC9,0xD6,0xE5),False)],
    [("높아진 변동성을 수확하는 오픈형 레버리지",15,RGBColor(0xC9,0xD6,0xE5),False)]],sp=1.25)
txt(s,Inches(0.9),Inches(6.4),Inches(9),Inches(0.4),[[("2026.07  |  내부 검토용 초안 — 대외 배포 금지",12,GRAY,False)]])

# ---- 2 상품 개요 ----
s=S(); bg(s,WHITE); header(s,"01  PRODUCT","상품 개요")
ov=[["구  분","내  용"],
    ["명  칭","미래에셋 변동성 하베스트 레버리지 — 오픈형 (가칭)"],
    ["구  조","오픈형(설정·해지 상시) · BM = 기초자산 150%"],
    ["기초자산","KOSPI200 / 삼성전자 / SK하이닉스 (시리즈 별도 설정)"],
    ["편입비","100% ~ 200% (기본 150%) · 상승 시 축소, 하락 시 확대"],
    ["유  형","상품1 룰베이스(사전 확정 공식) / 상품2 옵션이론 기반 모형"],
    ["리 셋","분기(63영업일)마다 기준 재설정 → 편입비 150% 복귀 (주간 분할 집행)"],
    ["상품 유형","국내 주식형 파생 (공모 / 사모 · 일임)"],
    ["가정 변동성","KOSPI200 연 50% · 삼성전자/SK하이닉스 연 60%"]]
table(s,Inches(0.75),Inches(1.5),Inches(9.3),ov,[Inches(1.7),Inches(7.6)],rh=Inches(0.46))
txt(s,Inches(0.75),Inches(6.35),Inches(9.4),Inches(0.6),[
    [("※ 레버리지 상품으로 원금 손실(0~100%)이 발생할 수 있으며 손실은 투자자에게 귀속됩니다. 예금자보호법 적용 대상이 아닙니다.",8.5,GRAY,False)]])
footer(s,2)

# ---- 3 왜 이 상품인가 ----
s=S(); bg(s,WHITE); header(s,"02  WHY","왜 오픈형 변동성 하베스트인가")
bullets(s,Inches(0.75),Inches(1.5),Inches(9.3),Inches(2.2),[
    ("①","레버리지 수요는 확인됐다 — 삼성전자·SK하이닉스 레버리지 ETF에 국내외 약 40조원(6/30). 그러나 기존 2배 ETF는 고가매수·저가매도 구조로 변동성에 '잠식'된다.",NAVY,True),
    ("②","본 상품은 정반대 — 오르면 줄이고(고가매도) 빠지면 늘린다(저가매수). 높아진 변동성이 비용이 아니라 수익의 원천.",ORANGE,True),
    ("③","같은 10% 변동에 레버리지 ETF는 +20%p 순응매매, 본 상품은 12~16%p 역방향 매매 — 몸집이 커질수록 시장 안정화에도 기여.",DGRAY,False),
])
pic(s,os.path.join(OIMG,"open_sensitivity.png"),Inches(1.7),Inches(3.55),Inches(7.4))
footer(s,3)

# ---- 4 방법론 ① 상품1 ----
s=S(); bg(s,WHITE); header(s,"03  METHOD","방법론 ① 상품1 — 룰베이스 편입비")
bullets(s,Inches(0.75),Inches(1.5),Inches(4.3),Inches(4.6),[
    ("■","편입비 = 150% − 1.25 × (지수 등락률)",NAVY,True),
    ("","· 리셋일 지수=100 기준, 캡 100~200%",DGRAY,False),
    ("","· 지수 ±10% → 편입비 ∓12.5%p",DGRAY,False),
    ("","· 지수 ±40%에서 캡(100/200%) 도달 — 현 고변동 국면(±40% 등락 가정)에 맞춘 설계",DGRAY,False),
    ("■","사전 확정 공식 → 투명·단순, 임의 판단 배제",ORANGE,True),
    ("■","레버리지 영역(150%±)에서는 지수가 오르면 편입비가 자연 하락 → 실제 매매는 공식 변화분의 약 절반 수준(시장충격 최소화)",DGRAY,False),
])
pic(s,os.path.join(OIMG,"open_w_k200.png"),Inches(5.15),Inches(1.6),Inches(5.2))
txt(s,Inches(5.15),Inches(6.5),Inches(5.2),Inches(0.35),[[("지수대별 편입비 곡선 (KOSPI200, ±40% 설계)",9,GRAY,False)]])
footer(s,4)

# ---- 5 방법론 ② 상품2 ----
s=S(); bg(s,WHITE); header(s,"03  METHOD","방법론 ② 상품2 — 옵션이론 기반 편입비")
bullets(s,Inches(0.75),Inches(1.5),Inches(4.3),Inches(4.6),[
    ("■","옵션이론(만기 0.25년 고정)으로 편입비를 산출",NAVY,True),
    ("","· 리셋일에 편입비가 정확히 150%가 되도록 기준가를 역산 (σ50% → 103.4, σ60% → 104.8)",DGRAY,False),
    ("","· 만기를 고정(rolling)해 시간이 지나도 편입비 곡선이 항상 같은 모양 유지",DGRAY,False),
    ("■","곡선이 완만한 S자 — 하락 초반엔 적극 매수, 극단에선 자동 감속. 캡(100~200%)이 구조적으로 자동 충족",ORANGE,True),
    ("■","변동성 수준(σ)에 따라 곡선이 자동 조정 — 단위형 변동성 하베스트와 동일 계보의 모형",DGRAY,False),
])
pic(s,os.path.join(OIMG,"open_w_stock.png"),Inches(5.15),Inches(1.6),Inches(5.2))
txt(s,Inches(5.15),Inches(6.5),Inches(5.2),Inches(0.35),[[("지수대별 편입비 곡선 (삼성전자·SK하이닉스, σ60%)",9,GRAY,False)]])
footer(s,5)

# ---- 6 방법론 ③ 리셋·실제 매매 ----
s=S(); bg(s,WHITE); header(s,"03  METHOD","방법론 ③ 오픈형 리셋과 실제 매매")
bullets(s,Inches(0.75),Inches(1.5),Inches(4.3),Inches(4.6),[
    ("■","분기(63영업일)마다 기준을 재설정 → 편입비 150% 복귀",NAVY,True),
    ("","· 주간 분할(13트랜치) 집행으로 일시 매매 없이 부드럽게 복귀. 설정·해지 자금도 분할 흡수",DGRAY,False),
    ("■","실제 순매수는 공식 변화분보다 작다",ORANGE,True),
    ("","· 지수 +10% 시: 목표는 −12.5%p지만 편입비가 스스로 −6.5%p 내려가 실제 매도는 약 −6%p",DGRAY,False),
    ("","· 지수 −10% 시: 실제 매수 약 +4%p — 하락 시 실탄 소모가 작음",DGRAY,False),
    ("","· 급락 극단(−30% 이하)에서는 캡 준수를 위해 오히려 매도 — '떨어지는 칼'을 무한정 잡지 않는 자기제어",DGRAY,False),
])
pic(s,os.path.join(OIMG,"open_trade_k200.png"),Inches(5.15),Inches(1.6),Inches(5.2))
txt(s,Inches(5.15),Inches(6.5),Inches(5.2),Inches(0.35),[[("지수 이동 시 실제 순매수 (레버리지 ETF와 비교)",9,GRAY,False)]])
footer(s,6)

# ---- 7 백테스트 ① KOSPI200 ----
s=S(); bg(s,WHITE); header(s,"04  BACKTEST","실측 백테스트 ① KOSPI200 (최근 1년)")
pic(s,os.path.join(OIMG,"open_bt_k200.png"),Inches(0.55),Inches(1.45),Inches(6.1))
bt=[["구  분","수익률","연변동성","MDD"],
    ["기초 1배","+190.0%","47.5%","−21.6%"],
    ["BM 150% 고정","+346.7%","71.2%","−31.7%"],
    ["상품1 룰베이스","+337.2%","55.0%","−21.6%"],
    ["상품2 옵션모형","+324.9%","55.8%","−22.5%"]]
table(s,Inches(6.85),Inches(1.7),Inches(3.7),bt,[Inches(1.45),Inches(0.85),Inches(0.75),Inches(0.65)],fs=9.5,rh=Inches(0.36),boldrows=(3,4))
bullets(s,Inches(6.85),Inches(3.85),Inches(3.7),Inches(2.6),[
    ("▶","추세 상승장이라 원수익률은 BM150이 최고",DGRAY,False),
    ("▶","두 상품은 MDD가 기초 1배 수준(−21.6%)으로 BM(−31.7%) 대비 하락 방어 우수",NAVY,True),
    ("▶","변동성도 15%p 이상 낮아 위험조정 성과 우위",ORANGE,True),
],sz=10.5,gap=1.25)
txt(s,Inches(0.55),Inches(6.6),Inches(9.8),Inches(0.35),[[("기간 2025-06-27~2026-07-10 · 분기 리셋 · 거래비용 15bps · 금리 2.5% · 가격수익률(배당 미반영) · 데이터: 주가 시계열(실거래일)",8.5,GRAY,False)]])
footer(s,7)

# ---- 8 백테스트 ② SK하이닉스 ----
s=S(); bg(s,WHITE); header(s,"04  BACKTEST","실측 백테스트 ② SK하이닉스 (최근 1년)")
pic(s,os.path.join(OIMG,"open_bt_hynix.png"),Inches(0.55),Inches(1.45),Inches(6.1))
bt=[["구  분","수익률","연변동성","MDD"],
    ["기초 1배","+662%","77.0%","−28.9%"],
    ["BM 150% 고정","+1,561%","115.5%","−42.0%"],
    ["상품1 룰베이스","+1,107%","87.2%","−28.9%"],
    ["상품2 옵션모형","+1,164%","89.8%","−28.9%"]]
table(s,Inches(6.85),Inches(1.7),Inches(3.7),bt,[Inches(1.45),Inches(0.85),Inches(0.75),Inches(0.65)],fs=9.5,rh=Inches(0.36),boldrows=(3,4))
bullets(s,Inches(6.85),Inches(3.85),Inches(3.7),Inches(2.6),[
    ("▶","기초 +662% 대비 상품 +1,100%대 — 상승 참여 충분",DGRAY,False),
    ("▶","MDD −28.9% = 기초와 동일. BM150(−42.0%)의 급락 취약점을 해소",NAVY,True),
    ("▶","7월 급락 구간에서 완충 효과 실측 확인",ORANGE,True),
],sz=10.5,gap=1.25)
txt(s,Inches(0.55),Inches(6.6),Inches(9.8),Inches(0.35),[[("기간 2025-06-25~2026-07-10 · 분기 리셋 · 거래비용 15bps · 금리 2.5% · 가격수익률(배당 미반영) · 데이터: 주가 시계열(실거래일)",8.5,GRAY,False)]])
footer(s,8)

# ---- 9 리스크 ----
s=S(); bg(s,WHITE); header(s,"05  RISK","리스크 및 유의사항")
bullets(s,Inches(0.75),Inches(1.55),Inches(9.3),Inches(4),[
    ("■","레버리지 위험: 편입비 100~200%의 레버리지 상품으로, 기초자산 하락 시 손실이 1배 대비 확대될 수 있습니다.",NAVY,True),
    ("■","하락 시 편입 확대 위험: 하락 구간에서 편입을 늘리는 구조 특성상, 하락이 장기화·심화되면 손실이 커질 수 있습니다(캡·자기제어로 완화하나 제거되지 않음).",DGRAY,False),
    ("■","추세장 상대 열위: 강한 일방향 상승장에서는 고정 150% 대비 수익률이 낮을 수 있습니다(백테스트에서 확인).",DGRAY,False),
    ("■","변동성 축소 위험: 시장 변동성이 낮아지면 변동성 매매 수익 기회가 줄어듭니다.",DGRAY,False),
    ("■","모형·운용 위험: 편입비 공식/모형의 가정(변동성 등)이 실제와 다를 경우 기대 성과와 차이가 발생할 수 있으며, 거래비용·유동성에 따라 결과가 달라질 수 있습니다.",DGRAY,False),
    ("■","백테스트 한계: 과거 성과는 미래 수익을 보장하지 않으며, 백테스트는 가정(비용·금리·배당 제외 등)에 따라 실제와 다를 수 있습니다.",DGRAY,False),
])
rect(s,Inches(0.75),Inches(5.75),Inches(9.3),Inches(0.95),LIGHT)
txt(s,Inches(1.0),Inches(5.86),Inches(8.9),Inches(0.75),[
    [("본 상품은 레버리지형 고위험 상품으로 원금의 전부(0~100%) 손실이 발생할 수 있으며, 그 손실은 투자자에게 귀속됩니다.",11.5,RED,True)],
    [("투자 전 상품설명서·약관·계약권유문서를 반드시 확인하시기 바랍니다.",10.5,GRAY,False)]],sp=1.2)
footer(s,9)

# ---- 10 운용부서/컴플 ----
s=S(); bg(s,NAVY)
rect(s,Inches(0.9),Inches(1.4),Inches(0.14),Inches(0.55),ORANGE)
txt(s,Inches(1.15),Inches(1.4),Inches(9),Inches(0.6),[[("운용부서 및 유의사항",22,WHITE,True)]])
txt(s,Inches(1.15),Inches(2.6),Inches(9),Inches(2),[
    [("AI금융공학운용부문 전략운용본부",15,CYAN,True)],
    [("Mission: 투자자 입장에서 장기 플러스 수익 달성",13,WHITE,False)],
    [("미래에셋 QPMS Vol Trading 시스템 기반 운용 · One Team Approach",12,RGBColor(0xC9,0xD6,0xE5),False)]],sp=1.4)
txt(s,Inches(1.15),Inches(5.3),Inches(8.9),Inches(1.6),[
    [("· 본 자료는 컴플라이언스 검토 전 내부 검토용 초안이며 대외 배포를 금합니다.",10,LORANGE,False)],
    [("· 펀드는 예금자보호법에 따라 보호되지 않으며, 자산가격 변동에 따라 원금손실(0~100%)이 발생할 수 있습니다.",10,RGBColor(0xC9,0xD6,0xE5),False)],
    [("· 상기 시뮬레이션·백테스트는 가정에 의거한 예시로 수익을 보장하지 않으며 실제 결과와 다를 수 있습니다.",10,RGBColor(0xC9,0xD6,0xE5),False)]],sp=1.3)
footer(s,10,"")

out=os.path.join(BASE,"변동성하베스트_오픈형_초안.pptx")
prs.save(out); print("saved:",out)
