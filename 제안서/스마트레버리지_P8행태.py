# -*- coding: utf-8 -*-
"""P8 재작성: 레버리지의 장점(추세 복리)은 개인 패턴과 정반대 — 설정(상장주식수) 실증 2컷
- 좌: 상승 추세장(2025.04~2026.06) — 오르면 매도(수익 확정) → 복리 못 누림
- 우: 급락·회복장(2020.01~2021.01) — 빠지면 원상복귀까지 버팀 → 잠식만 노출
데이터: ETF_데이터_pivot (***** 레버리지 ETF 상장주식수 · KOSPI200)
"""
import sys, io, os, platform
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
import numpy as np, pandas as pd
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
plt.rcParams['font.family']=('Malgun Gothic' if platform.system()=='Windows' else 'AppleGothic')
plt.rcParams['axes.unicode_minus']=False
NAVY='#043B72'; ORANGE='#F58220'; GRAY='#84888B'; AMBER='#8a4a08'
os.chdir(os.path.join("C:/Users/gamer38/Documents/Claude/Projects","변동성펀드","제안서"))
df=pd.read_parquet("C:/Users/gamer38/Documents/Claude/Projects/ETF WEB/ETF_데이터_pivot.parquet")
def px(nm,col='수정주가(원)',zero_ok=False):
    k=df[df['종목명']==nm].dropna(subset=[col]).sort_values('날짜')
    s=pd.Series(k[col].astype(float).values,index=pd.to_datetime(k['날짜']).values)
    i=pd.to_datetime(s.index); s=s[i.dayofweek<5]
    if zero_ok: return s
    chg=s.pct_change().fillna(1.0); return s[chg!=0]
k1=px('KODEX 200'); k2=px('KODEX 레버리지'); sh=px('KODEX 레버리지','상장주식수(주)',zero_ok=True)
RED='#C0392B'

def panel(ax,a,b):
    w=k1[(k1.index>=a)&(k1.index<=b)]; v=sh[(sh.index>=a)&(sh.index<=b)]
    lv=k2[(k2.index>=a)&(k2.index<=b)]
    ix=w/w.iloc[0]*100; lx=lv/lv.iloc[0]*100
    ax.fill_between(v.index,v.values/1e6,color='#f3c69a',alpha=0.9,label='***** 레버리지 ETF 상장주식수 (백만주, 좌)')
    ax.set_ylabel('상장주식수 (백만주)',fontsize=8.5,color=AMBER); ax.tick_params(axis='y',labelcolor=AMBER,labelsize=8.5)
    ax2=ax.twinx()
    ax2.plot(ix.index,ix.values,color=NAVY,lw=1.7,label='KOSPI200 (우, 시작=100)')
    ax2.set_ylabel('KOSPI200 (시작=100)',fontsize=8.5,color=NAVY); ax2.tick_params(axis='y',labelcolor=NAVY,labelsize=8.5)
    for spn in ['top']: ax.spines[spn].set_visible(False); ax2.spines[spn].set_visible(False)
    ax.grid(axis='y',color='#eee',lw=0.6); ax.set_axisbelow(True); ax.tick_params(axis='x',labelsize=8.5)
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%y.%m'))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
    return w,v,ix,lx,ax2

fig,(a1,b1)=plt.subplots(1,2,figsize=(10.4,3.5),dpi=150)

# 좌: 상승 추세장 2025-04-30 ~ 2026-06-22
a,b=pd.Timestamp('2025-04-30'),pd.Timestamp('2026-06-22')
w,v,ix,lx,ax2=panel(a1,a,b)
g=float(w.iloc[-1]/w.iloc[0]-1); gl2=float(lx.iloc[-1])-100
s0,s1=float(v.iloc[0])/1e6,float(v.iloc[-1])/1e6
a1.set_title(f"① 상승 추세장 : 오르면 팔아버린다\n(2025.04 ~ 2026.06 · KOSPI200 {g*100:+.0f}%)",fontsize=10.5,fontweight='bold',color='#222')
a1.annotate(f"주식수 {s0:,.0f} → {s1:,.0f}백만주 ({(s1/s0-1)*100:+.0f}%)\n= 수익 확정 매도 → 복리를 못 누린다",
            xy=(0.03,0.86),xycoords='axes fraction',fontsize=9.5,fontweight='bold',color=AMBER)
a1.set_xlim(a,b+pd.Timedelta(days=70))
ax2.annotate(f"지수\n{g*100:+.0f}%",xy=(ix.index[-1],float(ix.iloc[-1])),xytext=(5,-14),
             textcoords='offset points',fontsize=9.5,fontweight='bold',color=NAVY)
ymax=v.max()/1e6; a1.set_ylim(0,ymax*1.55)
print(f"[좌] 지수 {g*100:+.0f}% · 레버리지 {gl2:+.0f}% · 주식수 {s0:,.0f}→{s1:,.0f}백만주 ({(s1/s0-1)*100:+.0f}%)")

# 우: 급락·회복장 2020-01-02 ~ 2020-08-31
a,b=pd.Timestamp('2020-01-02'),pd.Timestamp('2020-08-31')
w,v,ix,lx,ax2=panel(b1,a,b)
lvl0=float(w.iloc[0]); tmin=w.idxmin()
rec=w[(w.index>tmin)&(w>=lvl0)]
trec=rec.index[0]
smax=float(v.max())/1e6; s0=float(v.iloc[0])/1e6; send=float(v.iloc[-1])/1e6
gmin=float(w.min()/lvl0-1); g=float(w.iloc[-1]/lvl0-1); gl2=float(lx.iloc[-1])-100
b1.set_title(f"② 급락·회복장 : 원상복귀까지 버틴다\n(2020.01 ~ 2020.08 · 급락 {gmin*100:+.0f}% 후 회복)",fontsize=10.5,fontweight='bold',color='#222')
ax2.axhline(100,color='#b8c6d4',lw=0.9,ls='--')
ax2.set_ylim(38,132)
ax2.annotate(f"원상복귀\n{trec.strftime('%Y.%m.%d')}",xy=(trec,100),xytext=(26,-74),textcoords='offset points',
             fontsize=8.5,fontweight='bold',color=NAVY,arrowprops=dict(arrowstyle='-',color=NAVY,lw=0.8))
b1.annotate(f"급락하자 주식수 {s0:,.0f} → {smax:,.0f}백만주 ({(smax/s0-1)*100:+.0f}%)\n= 원상복귀까지 버티기 → 잠식만 얻는다",
            xy=(0.26,0.88),xycoords='axes fraction',fontsize=9.5,fontweight='bold',color=AMBER)
ax2.annotate(f"지수 {g*100:+.1f}% (복귀·상회)",xy=(b-pd.Timedelta(days=4),113),ha='right',
             fontsize=9.5,fontweight='bold',color=NAVY)
b1.xaxis.set_major_locator(mdates.MonthLocator(interval=1))
ymax=v.max()/1e6; b1.set_ylim(0,ymax*1.55)
print(f"[우] 급락 {gmin*100:+.0f}% · 지수 {g*100:+.1f}% · 레버리지 {gl2:+.1f}% · 원상복귀 {trec.date()} · 주식수 {s0:,.0f}→최대 {smax:,.0f}→기말 {send:,.0f}백만주")

import matplotlib.lines as mlines
h1,l1=a1.get_legend_handles_labels()
fig.legend(h1+[mlines.Line2D([0],[0],color=NAVY,lw=1.7)],
           l1+['KOSPI200 (우, 시작=100)'],
           fontsize=8.5,loc='lower center',ncol=2,frameon=False,bbox_to_anchor=(0.5,-0.02))
plt.tight_layout(rect=[0,0.06,1,1])
plt.savefig(os.path.join('img','lev15_flow2.png')); plt.close()
print("차트 저장: lev15_flow2.png")

# ---- P8 재작성 ----
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
DST="Tech_TOP10_스마트액티브레버리지15_v1_20260812.pptx"
prs=Presentation(DST)
DARK=RGBColor(0x22,0x22,0x22); GRAY_=RGBColor(0x84,0x88,0x8B); ORANGE_=RGBColor(0xE8,0x72,0x0C)
FT_B="KoPub돋움체_Pro Bold"; FT_M="KoPub돋움체_Pro Medium"
SW=prs.slide_width/914400; SH=prs.slide_height/914400
def tbox(s,L,T,W,H,text,size,bold,color,font=None,sp=1.15):
    b=s.shapes.add_textbox(Inches(L),Inches(T),Inches(W),Inches(H))
    tf=b.text_frame; tf.word_wrap=True
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
        if font: r.font.name=font
        p.line_spacing=sp
    return b
s=prs.slides[7]
for sh_ in list(s.shapes): s.shapes._spTree.remove(sh_._element)
tbox(s,0.54,0.42,SW-1.1,0.45,"2. 그러나 레버리지의 장점은 개인의 매매 패턴과 맞지 않는다",21,True,DARK,FT_B)
tbox(s,0.55,0.98,SW-1.1,0.6,
     "① 상승하면 수익을 확정해 버려(매도) 추세장의 복리 효과를 누리지 못하고, ② 하락하면 지수가 원상복귀될 때까지 버팁니다. "
     "결국 복리는 놓치고 변동성 잠식만 얻습니다 — 레버리지 ETF의 실제 설정(상장주식수)으로 확인됩니다.",13.5,None,DARK)
s.shapes.add_picture(os.path.join('img','lev15_flow2.png'),Inches(0.45),Inches(1.98),width=Inches(9.95))
tbox(s,0.8,5.95,SW-1.7,0.5,"상승 때는 +1,437%를 두고 일찍 떠났고, 하락 때는 버텨서 잠식만 얻었습니다 — 장점과 정반대로 움직이는 것입니다.",13.5,True,ORANGE_,FT_B)
tbox(s,0.55,SH-0.62,SW-1.1,0.4,
     "※ ***** 레버리지 ETF 상장주식수·수정주가, ***** 200 ETF 수정주가 실측 (좌 2025.04.30~2026.06.22, 우 2020.01.02~2020.08.31) · 단순 2배 = 지수 누적수익률×2 (가상) · 참고: 월간 지수수익률 vs 주식수 증감률 상관 -0.56 (2019~2026) · 데이터: ETF데이터",8,False,GRAY_,FT_M)
prs.save(DST)
print("P8 저장 완료")
