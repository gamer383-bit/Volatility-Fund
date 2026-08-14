# -*- coding: utf-8 -*-
"""IT TOP10 스마트 액티브 레버리지 덱 — P5·6·7 교체 (다른 페이지는 건드리지 않음)"""
import sys, io, os
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
os.chdir(r"C:\Users\gamer38\Documents\Claude\Projects\변동성펀드\제안서")
DST="IT_TOP10_스마트액티브레버리지15_v1_20260812.pptx"
prs=Presentation(DST)
DARK=RGBColor(0x22,0x22,0x22); NAVY=RGBColor(0x04,0x3B,0x72); GRAY=RGBColor(0x84,0x88,0x8B)
ORANGE=RGBColor(0xE8,0x72,0x0C)
FT_B="KoPub돋움체_Pro Bold"; FT_M="KoPub돋움체_Pro Medium"
SW=prs.slide_width/914400; SH=prs.slide_height/914400
print(f"슬라이드 {SW:.2f}x{SH:.2f}")

def clear(s):
    for sh in list(s.shapes):
        s.shapes._spTree.remove(sh._element)
def tbox(s,L,T,W,H,text,size,bold,color,font=None,sp=1.15):
    b=s.shapes.add_textbox(Inches(L),Inches(T),Inches(W),Inches(H))
    tf=b.text_frame; tf.word_wrap=True
    for i,ln in enumerate(text.split('\n')):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.add_run(); r.text=ln
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
        if font: r.font.name=font
        p.line_spacing=sp
    return b
def page(idx,title,sub,imgs=None,foot=None,bigquote=None):
    s=prs.slides[idx-1]; clear(s)
    tbox(s,0.54,0.42,SW-1.1,0.45,title,21,True,DARK,FT_B)
    if sub: tbox(s,0.55,0.98,SW-1.1,0.6,sub,13.5,None,DARK)
    if imgs:
        for f,L,T,W in imgs:
            s.shapes.add_picture(os.path.join('img',f),Inches(L),Inches(T),width=Inches(W))
    if bigquote:
        tbox(s,0.8,SH-1.25,SW-1.7,0.5,bigquote,14,True,ORANGE,FT_B)
    tbox(s,0.55,SH-0.62,SW-1.1,0.4,foot or "※ 상기 자료는 이해를 돕기 위한 예시이며, 실제 투자성과와 다를 수 있습니다. 레버리지 상품은 원금 손실 위험이 큽니다.",8,False,GRAY,FT_M)

# ---- P5 : 긴 구간 변동성 잠식 (실측, 3분할) ----
g=0.14; w3=(SW-0.9-2*g)/3
page(5,"투자자가 아는 사실 : 지수가 제자리로 와도, 레버리지는 오지 않는다",
     "기간이 길수록 잠식은 커집니다 — ① KOSPI200이 4년 5개월 걸려 제자리로 왔을 때 레버리지는 -27.1%  "
     "② 하이닉스 레버리지는 출시 5주 뒤 주가가 +7.4% 올랐는데도 출시가로 회귀  ③ 왕복 5주 만에 -15.2%.",
     imgs=[('lev15_long_k200.png',0.45,1.95,w3),('lev15_hynix_a.png',0.45+w3+g,1.95,w3),('lev15_hynix_b.png',0.45+2*(w3+g),1.95,w3)],
     bigquote="잠식은 '왕복'에 청구되는 요금입니다 — 기간이 길수록, 변동이 클수록 요금은 커집니다.",
     foot="※ 수정주가 기준 실측 · KODEX 200/KODEX 레버리지, KODEX SK하이닉스단일종목레버리지 (2026.05.27 상장) · SK하이닉스 주가는 레버리지 ETF 일별수익률÷2로 역산 · 데이터: ETF데이터 · 상기 자료는 이해를 돕기 위한 예시입니다.")
print("P5 완료")

# ---- P6 : 변동성 잠식의 원리 (ETF웹 '레버리지 ETF' 레슨, 책 3-1 참조) ----
page(6,"투자자가 모르는 사실 ① : 방향이 없어도 원금이 깎인다 — 변동성 잠식",
     "레버리지는 '일별 수익률의 2배(1.5배)'를 맞추려 매일 오른 날 더 사고(고가 매수) 내린 날 팝니다(저가 매도). "
     "지수가 2배 갔다 제자리면 0이 되고, 하루 +10% 갔다 제자리만 와도 -1.8% — 이 손실이 등락마다 복리로 쌓입니다.",
     imgs=[('lev15_arith0.png',0.45,1.95,w3),('lev15_k200flat.png',0.45+w3+g,1.95,w3),('lev15_k200trend.png',0.45+2*(w3+g),1.95,w3)],
     bigquote="등락 반복에선 지수가 제자리여도 녹아내리고, 강하고 꾸준한 추세일 때만 진가 — 문제는 그런 추세가 드물다는 것입니다.",
     foot="※ KODEX 200·KODEX 레버리지 수정주가 실측 (제자리 사례 2026.02.26~04.16, 추세 사례 2026.03.31~05.29) · 단순 2배 = 기초 누적수익률×2 (가상) · '하루 +10% 후 제자리 = -1.8%'는 2배 레버리지 산술 예시 · 데이터: ETF데이터")
print("P6 완료")

# ---- P7 : 개인 습성 실증 (상장주식수 순증) ----
page(7,"투자자가 모르는 사실 ② : 개인의 습성은 레버리지와 정반대로 움직인다",
     "KODEX 레버리지 상장주식수(순증) 실증 — 오르면 팔고, 빠지면 사서 버팁니다. "
     "지수와 주식수 증감의 월간 상관 -0.56 (지수 +2%↑ 월 주식수 평균 -13.3%, -2%↓ 월 +18.8%).",
     imgs=[('lev15_flow.png',0.55,1.95,SW-1.1)],
     bigquote="잠식이 쌓이는 하락·왕복장엔 끝까지 남고, 복리가 쌓이는 추세 상승장(+342%)에선 주식수 -72% — 일찍 떠났습니다.",
     foot="※ KODEX 레버리지 상장주식수·KODEX 200 수정주가 (2019.01~2026.08) · 상관은 월말 기준 월간 지수수익률 vs 상장주식수 증감률 · 데이터: ETF데이터 · 상기 자료는 이해를 돕기 위한 예시입니다.")
print("P7 완료")

prs.save(DST)
print("저장:",DST)
